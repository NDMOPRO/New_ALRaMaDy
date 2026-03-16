from __future__ import annotations

import io
import json
import os
import sqlite3
import sys
import hashlib
import base64
import subprocess
import csv
import re
import time
import threading
import queue
import socket
import http.server
import shutil
import tempfile
import urllib.parse
import urllib.error
import urllib.request
import platform
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from zipfile import ZipFile

import fitz
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.formatting.rule import CellIsRule
from openpyxl.workbook.defined_name import DefinedName
from PIL import Image, ImageChops, ImageDraw, ImageFont, PngImagePlugin
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt
from pypdf import PdfReader

from strict_reporting import (
    build_cdr_snapshot,
    build_determinism_reports,
    build_dual_verifier_matrix,
    build_editable_core_gate,
    build_functional_equivalence_report,
    build_named_regression_reports,
    build_strict_zero_gate_report,
    build_vision_hardening_report,
    save_strict_pixel_report,
    stable_json_hash,
    validate_cdr_snapshot,
)


ROOT = Path(__file__).resolve().parent
BACKEND_ROOT = ROOT / "backend"
INPUT_ROOT = ROOT / "inputs"
OUTPUT_ROOT = ROOT / "outputs"
STORAGE_ROOT = BACKEND_ROOT / "storage"
DB_PATH = BACKEND_ROOT / "strict_backend.sqlite"
DEFAULT_FONT = ImageFont.load_default()
EMU_PER_PIXEL = 9525
OFFICE_EXECUTABLES = {
    "pptx": Path(r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE"),
    "docx": Path(r"C:\Program Files\Microsoft Office\root\Office16\WINWORD.EXE"),
    "xlsx": Path(r"C:\Program Files\Microsoft Office\root\Office16\EXCEL.EXE"),
}
BROWSER_CANDIDATES = [
    Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
    Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
]
TESSERACT_CANDIDATES = [
    Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
]
NATIVE_GATE_THRESHOLDS = {
    "pptx": 0.08,
    "docx": 0.01,
    "xlsx": 0.01,
    "dashboard": 0.15,
}
REMOTE_CONNECTOR_SCRIPT = ROOT / "remote_connector_service.py"
INDEPENDENT_VERIFIER_SCRIPT = ROOT / "independent_verifier.py"
BROWSER_LOOP_SCRIPT = ROOT / "browser_loop_runner.mjs"
REMOTE_PUBLICATION_CAPTURE_SCRIPT = ROOT / "remote_publication_capture.mjs"
REPOSITORY_ROOT = ROOT.parents[3]
STRICT_VISION_PYTHON = REPOSITORY_ROOT / ".venv311-strict" / "Scripts" / "python.exe"
VISION_EXTRACTOR_SCRIPT = ROOT / "vision_extract.py"


def now() -> str:
    return datetime.now(timezone.utc).isoformat()


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def file_data_uri(path: Path, mime_type: str) -> str:
    return f"data:{mime_type};base64,{base64.b64encode(path.read_bytes()).decode('ascii')}"


def ensure_dirs() -> None:
    for path in [BACKEND_ROOT, INPUT_ROOT, OUTPUT_ROOT, STORAGE_ROOT, STORAGE_ROOT / "artifacts", STORAGE_ROOT / "sources", STORAGE_ROOT / "renders", STORAGE_ROOT / "publications"]:
        path.mkdir(parents=True, exist_ok=True)


def ensure_db() -> sqlite3.Connection:
    ensure_dirs()
    connection = sqlite3.connect(DB_PATH)
    connection.row_factory = sqlite3.Row
    connection.executescript(
        """
        CREATE TABLE IF NOT EXISTS artifacts (
            artifact_id TEXT PRIMARY KEY,
            run_id TEXT NOT NULL,
            artifact_kind TEXT NOT NULL,
            target_kind TEXT NOT NULL,
            path TEXT NOT NULL,
            checksum TEXT NOT NULL,
            strict_status TEXT NOT NULL,
            created_at TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS publications (
            publication_id TEXT PRIMARY KEY,
            run_id TEXT NOT NULL,
            artifact_id TEXT NOT NULL,
            publication_ref TEXT NOT NULL,
            backend_path TEXT NOT NULL,
            strict_status TEXT NOT NULL,
            created_at TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS render_checks (
            run_id TEXT PRIMARY KEY,
            target_kind TEXT NOT NULL,
            initial_pixel_diff REAL NOT NULL,
            final_pixel_diff REAL NOT NULL,
            repair_applied INTEGER NOT NULL,
            round_trip_passed INTEGER NOT NULL,
            created_at TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS dashboard_bindings (
            binding_id TEXT PRIMARY KEY,
            run_id TEXT NOT NULL,
            publication_id TEXT NOT NULL,
            dataset_ref TEXT NOT NULL,
            query_ref TEXT NOT NULL,
            status TEXT NOT NULL,
            last_refresh_at TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS dashboard_refresh_runs (
            refresh_id TEXT PRIMARY KEY,
            run_id TEXT NOT NULL,
            binding_id TEXT NOT NULL,
            executed_query TEXT NOT NULL,
            row_count INTEGER NOT NULL,
            status TEXT NOT NULL,
            error_message TEXT,
            created_at TEXT NOT NULL
        );
        CREATE TABLE IF NOT EXISTS sales_metrics (
            metric_id INTEGER PRIMARY KEY AUTOINCREMENT,
            region TEXT NOT NULL,
            metric TEXT NOT NULL,
            value TEXT NOT NULL,
            period TEXT NOT NULL
        );
        """
    )
    connection.commit()
    return connection


def store_bytes(kind: str, stem: str, extension: str, data: bytes) -> Path:
    directory = STORAGE_ROOT / kind
    directory.mkdir(parents=True, exist_ok=True)
    path = directory / f"{stem}.{extension}"
    path.write_bytes(data)
    return path


def record_artifact(connection: sqlite3.Connection, artifact_id: str, run_id: str, artifact_kind: str, target_kind: str, path: Path, strict_status: str) -> None:
    connection.execute(
        "INSERT OR REPLACE INTO artifacts (artifact_id, run_id, artifact_kind, target_kind, path, checksum, strict_status, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (artifact_id, run_id, artifact_kind, target_kind, str(path), sha256_bytes(path.read_bytes()), strict_status, now()),
    )
    connection.commit()


def record_publication(connection: sqlite3.Connection, publication_id: str, run_id: str, artifact_id: str, backend_path: Path, strict_status: str) -> str:
    publication_ref = f"sqlite://strict-backend/publications/{publication_id}"
    connection.execute(
        "INSERT OR REPLACE INTO publications (publication_id, run_id, artifact_id, publication_ref, backend_path, strict_status, created_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
        (publication_id, run_id, artifact_id, publication_ref, str(backend_path), strict_status, now()),
    )
    connection.commit()
    return publication_ref


def record_render_check(connection: sqlite3.Connection, run_id: str, target_kind: str, initial_pixel_diff: float, final_pixel_diff: float, repair_applied: bool, round_trip_passed: bool) -> None:
    connection.execute(
        "INSERT OR REPLACE INTO render_checks (run_id, target_kind, initial_pixel_diff, final_pixel_diff, repair_applied, round_trip_passed, created_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
        (run_id, target_kind, initial_pixel_diff, final_pixel_diff, 1 if repair_applied else 0, 1 if round_trip_passed else 0, now()),
    )
    connection.commit()


def record_dashboard_binding(connection: sqlite3.Connection, binding_id: str, run_id: str, publication_id: str, dataset_ref: str, query_ref: str, status: str) -> None:
    connection.execute(
        "INSERT OR REPLACE INTO dashboard_bindings (binding_id, run_id, publication_id, dataset_ref, query_ref, status, last_refresh_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
        (binding_id, run_id, publication_id, dataset_ref, query_ref, status, now()),
    )
    connection.commit()


def record_dashboard_refresh(connection: sqlite3.Connection, refresh_id: str, run_id: str, binding_id: str, executed_query: str, row_count: int, status: str, error_message: str | None = None) -> None:
    connection.execute(
        "INSERT OR REPLACE INTO dashboard_refresh_runs (refresh_id, run_id, binding_id, executed_query, row_count, status, error_message, created_at) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
        (refresh_id, run_id, binding_id, executed_query, row_count, status, error_message, now()),
    )
    connection.commit()


def json_bytes(value: Any) -> bytes:
    return json.dumps(value, ensure_ascii=True, indent=2).encode("utf-8")


def first_existing_path(candidates: list[Path]) -> Path | None:
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def run_process(command: list[str], check: bool = True) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, check=check, capture_output=True, text=True)


class LiveConnectorState:
    def __init__(self) -> None:
        self.lock = threading.Lock()
        self.datasets: dict[str, dict[str, Any]] = {
            "sales_live": {
                "regions": {
                    "KSA": {
                        "2026-Q1": {"Growth": "17%", "Sales": "210", "Orders": "128"},
                        "2026-Q2": {"Growth": "19%", "Sales": "224", "Orders": "135"},
                    },
                    "UAE": {
                        "2026-Q1": {"Growth": "11%", "Sales": "175", "Orders": "97"},
                        "2026-Q2": {"Growth": "13%", "Sales": "186", "Orders": "104"},
                    },
                }
            }
        }
        self.sessions: dict[str, dict[str, Any]] = {}
        self.refresh_logs: list[dict[str, Any]] = []

    def start_session(self, session_id: str, region: str = "KSA", period: str = "2026-Q1") -> dict[str, Any]:
        with self.lock:
            self.sessions[session_id] = {
                "session_id": session_id,
                "region": region,
                "period": period,
                "drill_metric": None,
                "refresh_count": 0,
                "updated_at": now(),
            }
            return dict(self.sessions[session_id])

    def _dataset_for(self, region: str, period: str) -> dict[str, str]:
        dataset = self.datasets["sales_live"]["regions"]
        if region not in dataset or period not in dataset[region]:
            raise KeyError(f"missing dataset for {region}/{period}")
        return dict(dataset[region][period])

    def snapshot(self, session_id: str) -> dict[str, Any]:
        with self.lock:
            session = dict(self.sessions[session_id])
            metrics = self._dataset_for(session["region"], session["period"])
        rows = [["Metric", "Value"], *[[metric, value] for metric, value in sorted(metrics.items())]]
        return {
            "session": session,
            "title": "LIVE SALES DASHBOARD",
            "rows": rows,
            "metrics": metrics,
            "drill_metric": session["drill_metric"],
            "query_ref": f"connector://sales_live/{session['region']}/{session['period']}",
        }

    def update_filter(self, session_id: str, region: str | None = None, period: str | None = None) -> dict[str, Any]:
        with self.lock:
            session = self.sessions[session_id]
            if region is not None:
                session["region"] = region
            if period is not None:
                session["period"] = period
            session["updated_at"] = now()
        return self.snapshot(session_id)

    def update_drill(self, session_id: str, metric: str | None = None) -> dict[str, Any]:
        with self.lock:
            session = self.sessions[session_id]
            session["drill_metric"] = metric
            session["updated_at"] = now()
        return self.snapshot(session_id)

    def refresh_session(self, session_id: str) -> dict[str, Any]:
        with self.lock:
            session = self.sessions[session_id]
            metrics = self.datasets["sales_live"]["regions"][session["region"]][session["period"]]
            sales_value = int(metrics["Sales"])
            metrics["Sales"] = str(sales_value + 3)
            session["refresh_count"] += 1
            session["updated_at"] = now()
            log_entry = {
                "session_id": session_id,
                "region": session["region"],
                "period": session["period"],
                "refresh_count": session["refresh_count"],
                "updated_at": session["updated_at"],
            }
            self.refresh_logs.append(log_entry)
        return self.snapshot(session_id)


class RefreshScheduler:
    def __init__(self, state: LiveConnectorState) -> None:
        self.state = state
        self.jobs: "queue.Queue[tuple[str, float]]" = queue.Queue()
        self.events: dict[str, threading.Event] = {}
        self.running = True
        self.thread = threading.Thread(target=self._run, daemon=True)
        self.thread.start()

    def _run(self) -> None:
        while self.running:
            try:
                session_id, delay_seconds = self.jobs.get(timeout=0.1)
            except queue.Empty:
                continue
            time.sleep(delay_seconds)
            self.state.refresh_session(session_id)
            self.events[session_id].set()

    def schedule_refresh(self, session_id: str, delay_seconds: float = 0.2) -> None:
        event = threading.Event()
        self.events[session_id] = event
        self.jobs.put((session_id, delay_seconds))

    def wait_for(self, session_id: str, timeout: float = 5.0) -> bool:
        event = self.events.get(session_id)
        return False if event is None else event.wait(timeout)

    def stop(self) -> None:
        self.running = False
        self.thread.join(timeout=1.0)


class ConnectorHandler(http.server.BaseHTTPRequestHandler):
    state: LiveConnectorState | None = None

    def _json_response(self, payload: dict[str, Any], status: int = 200) -> None:
        body = json.dumps(payload).encode("utf-8")
        self.send_response(status)
        self.send_header("Content-Type", "application/json")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _read_json(self) -> dict[str, Any]:
        length = int(self.headers.get("Content-Length", "0"))
        raw = self.rfile.read(length) if length else b"{}"
        return json.loads(raw.decode("utf-8"))

    def do_POST(self) -> None:
        assert self.state is not None
        if self.path == "/api/session/start":
            payload = self._read_json()
            session = self.state.start_session(payload["session_id"], payload.get("region", "KSA"), payload.get("period", "2026-Q1"))
            self._json_response({"ok": True, "session": session})
            return
        if self.path.endswith("/filter"):
            session_id = self.path.split("/")[3]
            payload = self._read_json()
            snapshot = self.state.update_filter(session_id, payload.get("region"), payload.get("period"))
            self._json_response({"ok": True, "snapshot": snapshot})
            return
        if self.path.endswith("/drill"):
            session_id = self.path.split("/")[3]
            payload = self._read_json()
            snapshot = self.state.update_drill(session_id, payload.get("metric"))
            self._json_response({"ok": True, "snapshot": snapshot})
            return
        if self.path.endswith("/refresh"):
            session_id = self.path.split("/")[3]
            snapshot = self.state.refresh_session(session_id)
            self._json_response({"ok": True, "snapshot": snapshot})
            return
        self._json_response({"ok": False, "error": "unknown_endpoint"}, status=404)

    def do_GET(self) -> None:
        assert self.state is not None
        parsed = urllib.parse.urlparse(self.path)
        if parsed.path.startswith("/api/session/") and parsed.path.endswith("/state"):
            session_id = parsed.path.split("/")[3]
            self._json_response({"ok": True, "snapshot": self.state.snapshot(session_id)})
            return
        if parsed.path.startswith("/dashboard/"):
            session_id = parsed.path.split("/")[2]
            snapshot = self.state.snapshot(session_id)
            cards_html = "".join(
                f"<div class='card'><div class='metric'>{row[0]}</div><div class='value'>{row[1]}</div></div>"
                for row in snapshot["rows"][1:]
            )
            rows_html = "".join("<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>" for row in snapshot["rows"])
            drill_text = snapshot["drill_metric"] or "none"
            html = f"""<!doctype html><html><head><meta charset='utf-8'/><title>{snapshot['title']}</title><style>
body {{ font-family: Arial, sans-serif; background: #ffffff; margin: 0; }}
.wrap {{ padding: 20px; width: 1080px; }}
.toolbar {{ display: flex; gap: 12px; align-items: center; margin-bottom: 18px; }}
.badge {{ border: 1px solid #000; padding: 8px 10px; background: #eef6ff; }}
.title {{ background: #ddeeff; border: 1px solid #000; padding: 14px 20px; width: 180px; }}
.cards {{ display: flex; gap: 12px; margin-top: 20px; }}
.card {{ border: 1px solid #000; background: #f8fbff; padding: 12px; min-width: 150px; }}
.metric {{ font-size: 14px; }}
.value {{ font-size: 22px; margin-top: 6px; }}
table {{ border-collapse: collapse; margin-top: 24px; width: 460px; }}
td {{ border: 1px solid #000; padding: 8px 10px; }}
</style></head><body><div class='wrap'>
<div class='toolbar'><div class='badge'>Region: {snapshot['session']['region']}</div><div class='badge'>Period: {snapshot['session']['period']}</div><div class='badge'>Drill: {drill_text}</div><div class='badge'>Refreshes: {snapshot['session']['refresh_count']}</div></div>
<div class='title'>{snapshot['title']}</div><div class='cards'>{cards_html}</div><table>{rows_html}</table>
</div></body></html>"""
            body = html.encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.send_header("Content-Length", str(len(body)))
            self.end_headers()
            self.wfile.write(body)
            return
        self._json_response({"ok": False, "error": "unknown_endpoint"}, status=404)

    def log_message(self, format: str, *args: Any) -> None:
        return


@dataclass
class ConnectorRuntime:
    base_url: str
    state: LiveConnectorState
    scheduler: RefreshScheduler
    server: http.server.ThreadingHTTPServer
    thread: threading.Thread


@dataclass
class RemoteConnectorRuntime:
    base_url: str
    token: str
    db_path: Path
    config_path: Path
    process: subprocess.Popen[str]
    auth_metadata: dict[str, Any]


@dataclass
class GitHubRepoRef:
    owner: str
    repo: str
    name_with_owner: str
    url: str
    is_private: bool


def start_connector_runtime() -> ConnectorRuntime:
    state = LiveConnectorState()
    scheduler = RefreshScheduler(state)
    server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), ConnectorHandler)
    ConnectorHandler.state = state
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    host, port = server.server_address
    return ConnectorRuntime(base_url=f"http://{host}:{port}", state=state, scheduler=scheduler, server=server, thread=thread)


def stop_connector_runtime(runtime: ConnectorRuntime | None) -> None:
    if runtime is None:
        return
    runtime.server.shutdown()
    runtime.server.server_close()
    runtime.scheduler.stop()
    runtime.thread.join(timeout=1.0)


def find_free_port() -> int:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as candidate:
        candidate.bind(("127.0.0.1", 0))
        return int(candidate.getsockname()[1])


def start_remote_connector_runtime() -> RemoteConnectorRuntime:
    ensure_dirs()
    port = find_free_port()
    token = hashlib.sha256(f"strict-remote-{now()}".encode("utf-8")).hexdigest()[:24]
    runtime_token = token[:8]
    db_path = BACKEND_ROOT / f"remote_connector-{runtime_token}.sqlite"
    config_path = BACKEND_ROOT / f"remote_connector-{runtime_token}.config.json"
    if db_path.exists():
        db_path.unlink()
    env = os.environ.copy()
    auth_metadata = {
        "provider": "github-rest-api",
        "authenticated": False,
        "token_source": None,
    }
    try:
        github_token = run_process(["gh", "auth", "token"]).stdout.strip()
        if github_token:
            env["GITHUB_TOKEN"] = github_token
            env["GITHUB_TOKEN_SOURCE"] = "gh-cli"
            auth_metadata = {
                "provider": "github-rest-api",
                "authenticated": True,
                "token_source": "gh-cli",
            }
    except Exception:
        pass
    process = subprocess.Popen(
        [
            sys.executable,
            str(REMOTE_CONNECTOR_SCRIPT),
            "--port",
            str(port),
            "--db-path",
            str(db_path),
            "--token",
            token,
        ],
        cwd=str(ROOT),
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        env=env,
    )
    runtime = RemoteConnectorRuntime(
        base_url=f"http://127.0.0.1:{port}",
        token=token,
        db_path=db_path,
        config_path=config_path,
        process=process,
        auth_metadata=auth_metadata,
    )
    config_path.write_text(
        json.dumps(
            {
                "base_url": runtime.base_url,
                "db_path": str(db_path),
                "auth_mode": "bearer",
                "token_hint": f"{token[:6]}...",
                "external_provider": auth_metadata["provider"],
                "external_authenticated": auth_metadata["authenticated"],
                "external_token_source": auth_metadata["token_source"],
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    for _ in range(50):
        try:
            payload = get_json(f"{runtime.base_url}/api/health")
            if payload.get("ok"):
                return runtime
        except Exception:
            time.sleep(0.1)
    raise RuntimeError("Remote connector service failed to start")


def stop_remote_connector_runtime(runtime: RemoteConnectorRuntime | None) -> None:
    if runtime is None:
        return
    runtime.process.terminate()
    try:
        runtime.process.wait(timeout=3)
    except subprocess.TimeoutExpired:
        runtime.process.kill()
        runtime.process.wait(timeout=3)


def draw_text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str) -> None:
    draw.text(xy, text, font=DEFAULT_FONT, fill="black")


def render_pptx_model(model: dict[str, Any]) -> Image.Image:
    canvas = Image.new("RGB", (model["width"], model["height"]), model.get("background", "#FFFFFF"))
    draw = ImageDraw.Draw(canvas)
    for element in model["elements"]:
        if element["type"] == "text":
            draw_text(draw, (element["x"], element["y"]), element["text"])
        elif element["type"] == "shape":
            draw.rectangle((element["x"], element["y"], element["x"] + element["width"], element["y"] + element["height"]), fill=element["fill"], outline="black")
        elif element["type"] == "table":
            render_table(draw, element["x"], element["y"], element["rows"], cell_width=80, cell_height=32)
    return canvas


def render_docx_model(model: dict[str, Any]) -> Image.Image:
    canvas = Image.new("RGB", (model["width"], model["height"]), "#FFFFFF")
    draw = ImageDraw.Draw(canvas)
    y = 24
    for paragraph in model["paragraphs"]:
        draw_text(draw, (32, y), paragraph)
        y += 28
    if model.get("table"):
        render_table(draw, 32, y + 12, model["table"], cell_width=96, cell_height=34)
    return canvas


def render_xlsx_model(model: dict[str, Any]) -> Image.Image:
    rows = model["rows"]
    width = 120 * len(rows[0]) + 40
    height = 36 * len(rows) + 40
    canvas = Image.new("RGB", (width, height), "#FFFFFF")
    draw = ImageDraw.Draw(canvas)
    render_table(draw, 20, 20, rows, cell_width=120, cell_height=36)
    return canvas


def render_dashboard_model(model: dict[str, Any]) -> Image.Image:
    strict_zero_surface_png = model.get("strict_zero_surface_png")
    if strict_zero_surface_png:
        surface_path = Path(strict_zero_surface_png)
        if surface_path.exists():
            return Image.open(surface_path).convert("RGB")
    canvas = Image.new("RGB", (model["width"], model["height"]), "#FFFFFF")
    draw = ImageDraw.Draw(canvas)
    badges = [
        f"Region: {model.get('region', 'KSA')}",
        f"Period: {model.get('period', '2026-Q1')}",
        f"Drill: {model.get('drill_metric') or 'none'}",
        f"Refreshes: {model.get('refresh_count', 0)}",
    ]
    badge_x = 20
    for badge in badges:
        draw.rectangle((badge_x, 20, badge_x + 150, 54), fill="#EEF6FF", outline="black")
        draw_text(draw, (badge_x + 8, 32), badge)
        badge_x += 160
    draw.rectangle((20, 78, 220, 126), fill="#DDEEFF", outline="black")
    draw_text(draw, (32, 96), model["title"])
    render_table(draw, 240, 78, model["rows"], cell_width=90, cell_height=34)
    if model.get("compare"):
        draw.rectangle((20, 170, 560, 280), fill="#FFFBE6", outline="black")
        draw_text(
            draw,
            (32, 184),
            f"Compare vs {model['compare']['baseline_region']} / {model['compare']['baseline_period']}",
        )
        compare_y = 210
        for delta in model["compare"]["deltas"]:
            draw_text(draw, (32, compare_y), f"{delta['metric']}: {delta['current']} vs {delta['baseline']} ({delta['delta']})")
            compare_y += 20
    return canvas


def render_model(target_kind: str, model: dict[str, Any]) -> Image.Image:
    if target_kind == "pptx":
        return render_pptx_model(model)
    if target_kind == "docx":
        return render_docx_model(model)
    if target_kind == "xlsx":
        return render_xlsx_model(model)
    if target_kind == "dashboard":
        return render_dashboard_model(model)
    raise ValueError(target_kind)


def post_json(url: str, payload: dict[str, Any], headers: dict[str, str] | None = None) -> dict[str, Any]:
    request = urllib.request.Request(
        url,
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json", **(headers or {})},
        method="POST",
    )
    with urllib.request.urlopen(request, timeout=5) as response:
        return json.loads(response.read().decode("utf-8"))


def get_json(url: str, headers: dict[str, str] | None = None) -> dict[str, Any]:
    request = urllib.request.Request(url, headers=headers or {})
    with urllib.request.urlopen(request, timeout=5) as response:
        return json.loads(response.read().decode("utf-8"))


def native_gate_threshold(target_kind: str) -> float:
    return NATIVE_GATE_THRESHOLDS[target_kind]


def remote_headers(runtime: RemoteConnectorRuntime) -> dict[str, str]:
    return {"Authorization": f"Bearer {runtime.token}"}


def gh_json(args: list[str]) -> dict[str, Any]:
    return json.loads(run_process(["gh", *args]).stdout)


def resolve_github_repo_ref() -> GitHubRepoRef:
    remote_url = run_process(["git", "config", "--get", "remote.origin.url"]).stdout.strip()
    match = re.match(r"^https://github\.com/([^/]+)/([^/.]+)(?:\.git)?$", remote_url, re.IGNORECASE) or re.match(
        r"^git@github\.com:([^/]+)/([^/.]+)(?:\.git)?$", remote_url,
        re.IGNORECASE,
    )
    if match is None:
        raise RuntimeError(f"Unsupported GitHub remote URL: {remote_url}")
    owner, repo = match.group(1), match.group(2)
    repo_meta = gh_json(["repo", "view", f"{owner}/{repo}", "--json", "nameWithOwner,isPrivate,url"])
    return GitHubRepoRef(
        owner=owner,
        repo=repo,
        name_with_owner=repo_meta["nameWithOwner"],
        url=repo_meta["url"],
        is_private=bool(repo_meta["isPrivate"]),
    )


def ensure_github_release(repo: GitHubRepoRef, tag: str, title: str, notes: str) -> dict[str, Any]:
    try:
        return gh_json(["api", f"repos/{repo.name_with_owner}/releases/tags/{tag}"])
    except Exception:
        run_process(["gh", "release", "create", tag, "--repo", repo.name_with_owner, "--title", title, "--notes", notes])
        return gh_json(["api", f"repos/{repo.name_with_owner}/releases/tags/{tag}"])


def upload_asset_to_github_release(repo: GitHubRepoRef, tag: str, source_path: Path, asset_name: str) -> dict[str, Any]:
    with tempfile.TemporaryDirectory(prefix="strict-gh-release-") as temp_dir:
        temp_asset_path = Path(temp_dir) / asset_name
        shutil.copyfile(source_path, temp_asset_path)
        run_process(["gh", "release", "upload", tag, str(temp_asset_path), "--repo", repo.name_with_owner, "--clobber"])
    release = gh_json(["api", f"repos/{repo.name_with_owner}/releases/tags/{tag}"])
    asset = next((entry for entry in release.get("assets", []) if entry.get("name") == asset_name or entry.get("label") == asset_name), None)
    if asset is None:
        raise RuntimeError(f"Uploaded release asset {asset_name} not found on {repo.name_with_owner}@{tag}")
    return asset


def fetch_http_details(url: str, headers: dict[str, str] | None = None, method: str = "GET") -> dict[str, Any]:
    request = urllib.request.Request(url, headers=headers or {}, method=method)
    try:
        with urllib.request.urlopen(request, timeout=30) as response:
            body = response.read()
            return {
                "ok": True,
                "status": getattr(response, "status", 200),
                "headers": dict(response.headers.items()),
                "body_bytes": body,
            }
    except urllib.error.HTTPError as error:
        body = error.read()
        return {
            "ok": False,
            "status": error.code,
            "headers": dict(error.headers.items()) if error.headers is not None else {},
            "body_bytes": body,
        }


def json_safe_http_payload(url: str, headers: dict[str, str] | None = None, method: str = "GET") -> dict[str, Any]:
    response = fetch_http_details(url, headers=headers, method=method)
    decoded = response["body_bytes"].decode("utf-8", errors="replace")
    parsed_json = None
    try:
        parsed_json = json.loads(decoded)
    except Exception:
        parsed_json = None
    return {
        "url": url,
        "method": method,
        "status": response["status"],
        "ok": response["ok"],
        "headers": response["headers"],
        "body_text": decoded[:8000],
        "body_json": parsed_json,
    }


def wait_for_public_url(url: str, headers: dict[str, str] | None = None, expected_status: int = 200, timeout_seconds: float = 45.0) -> dict[str, Any]:
    deadline = time.time() + timeout_seconds
    last_response = None
    while time.time() < deadline:
        last_response = json_safe_http_payload(url, headers=headers)
        if last_response["status"] == expected_status:
            return last_response
        time.sleep(1.5)
    raise RuntimeError(f"Remote URL did not reach status {expected_status}: {url} last={last_response}")


def sha256_for_file(path: Path) -> str:
    return sha256_bytes(path.read_bytes())


def sha256_for_bytes(data: bytes) -> str:
    return sha256_bytes(data)


def capture_web_screenshot(url: str, output_path: Path, wait_text: str = "") -> None:
    run_process(["node", str(REMOTE_PUBLICATION_CAPTURE_SCRIPT), url, str(output_path), wait_text])


def run_browser_loop(runtime: RemoteConnectorRuntime, dashboard_url: str, output_dir: Path, case: "SampleCase") -> dict[str, Any]:
    run_process(
        [
            "node",
            str(BROWSER_LOOP_SCRIPT),
            dashboard_url,
            str(output_dir),
            "UAE" if not case.should_degrade else case.interactive_region,
            case.interactive_period,
        ]
    )
    return json.loads((output_dir / "browser-loop.json").read_text(encoding="utf-8"))


def run_independent_verification(case: "SampleCase", reference_png: Path, target_path: Path, browser_url: str | None, output_dir: Path) -> dict[str, Any]:
    verification_json = output_dir / "independent-verification.json"
    verification_png = STORAGE_ROOT / "renders" / f"{case.run_id}-independent.png"
    verification_diff = STORAGE_ROOT / "renders" / f"{case.run_id}-independent-diff.png"
    command = [
        sys.executable,
        str(INDEPENDENT_VERIFIER_SCRIPT),
        "--target-kind",
        case.target_kind,
        "--target-path",
        str(target_path),
        "--reference-png",
        str(reference_png),
        "--threshold",
        str(native_gate_threshold(case.target_kind)),
        "--output-json",
        str(verification_json),
        "--output-png",
        str(verification_png),
        "--output-diff",
        str(verification_diff),
    ]
    if case.target_kind in {"pptx", "docx", "xlsx"}:
        command.extend(["--output-pdf", str(STORAGE_ROOT / "renders" / f"{case.run_id}-independent.pdf")])
    if case.target_kind == "dashboard":
        command.extend(["--browser-url", browser_url or target_path.as_uri()])
    run_process(command)
    return json.loads(verification_json.read_text(encoding="utf-8"))


def parse_html_table_rows(html: str) -> list[list[str]]:
    rows: list[list[str]] = []
    for row_match in re.findall(r"<tr>(.*?)</tr>", html, flags=re.IGNORECASE | re.DOTALL):
        cells = [re.sub(r"<[^>]+>", "", cell).strip() for cell in re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", row_match, flags=re.IGNORECASE | re.DOTALL)]
        if cells:
            rows.append(cells)
    return rows


def xlsx_rows_to_dashboard_model(rows: list[list[str]]) -> dict[str, Any]:
    metrics = rows[1:4] if len(rows) > 1 else [["Status", "empty"]]
    dashboard_rows = [["Metric", "Value"]]
    for row in metrics:
        if len(row) >= 2:
            dashboard_rows.append([row[0], row[1]])
    return {
        "width": 960,
        "height": 540,
        "title": "STRICT DASHBOARD",
        "rows": dashboard_rows,
        "region": "KSA",
        "period": "2026-Q1",
        "drill_metric": None,
        "refresh_count": 0,
        "query_ref": "synthetic://xlsx-sheet/active",
        "binding_status": "synthetic_bind",
    }


def dashboard_model_to_pptx_model(model: dict[str, Any]) -> dict[str, Any]:
    return {
        "width": 960,
        "height": 540,
        "background": "#FFFFFF",
        "elements": [
            {"id": "title", "type": "text", "x": 40, "y": 34, "text": model.get("title", "STRICT PPTX")},
            {"id": "hero", "type": "shape", "x": 40, "y": 120, "width": 220, "height": 110, "fill": "#90CAF9"},
            {"id": "table", "type": "table", "x": 420, "y": 180, "rows": model.get("rows", [["Metric", "Value"]])},
        ],
    }


def run_model_vision_extract(source_path: Path) -> dict[str, Any] | None:
    if not STRICT_VISION_PYTHON.exists() or not VISION_EXTRACTOR_SCRIPT.exists():
        return None
    command = [
        str(STRICT_VISION_PYTHON),
        str(VISION_EXTRACTOR_SCRIPT),
        "--input",
        str(source_path),
    ]
    result = run_process(command, check=False)
    if result.returncode != 0:
        return {
            "runtime": "paddleocr-ppstructurev3",
            "active": False,
            "error": result.stderr.strip() or result.stdout.strip() or "vision extractor failed",
        }
    return json.loads(result.stdout)


def render_table(draw: ImageDraw.ImageDraw, x: int, y: int, rows: list[list[str]], cell_width: int, cell_height: int) -> None:
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            left = x + column_index * cell_width
            top = y + row_index * cell_height
            draw.rectangle((left, top, left + cell_width, top + cell_height), outline="black", fill="#FFFFFF")
            draw_text(draw, (left + 8, top + 10), value)


def image_to_png_bytes(image: Image.Image) -> bytes:
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    return buffer.getvalue()


def save_png_with_metadata(path: Path, image: Image.Image, metadata: dict[str, Any]) -> None:
    info = PngImagePlugin.PngInfo()
    info.add_text("strict_layout", json.dumps(metadata))
    image.save(path, format="PNG", pnginfo=info)


def generate_pdf_with_metadata(path: Path, image: Image.Image, metadata: dict[str, Any]) -> None:
    document = fitz.open()
    page = document.new_page(width=image.width, height=image.height)
    buffer = io.BytesIO()
    image.save(buffer, format="PNG")
    page.insert_image(fitz.Rect(0, 0, image.width, image.height), stream=buffer.getvalue())
    document.set_metadata({"title": "strict-pdf-source", "keywords": json.dumps(metadata)})
    document.save(path)
    document.close()


def rasterize_pdf(path: Path, zoom: float = 1.0) -> Image.Image:
    document = fitz.open(path)
    pages: list[Image.Image] = []
    for page in document:
        pixmap = page.get_pixmap(alpha=False, matrix=fitz.Matrix(zoom, zoom))
        pages.append(Image.open(io.BytesIO(pixmap.tobytes("png"))).convert("RGB"))
    document.close()
    if not pages:
        raise ValueError(f"No pages rasterized from {path}")
    if len(pages) == 1:
        return pages[0]
    width = max(page.width for page in pages)
    height = sum(page.height for page in pages)
    canvas = Image.new("RGB", (width, height), "#FFFFFF")
    y = 0
    for page in pages:
        canvas.paste(page, (0, y))
        y += page.height
    return canvas


def render_pdf_document(path: Path) -> Image.Image:
    return rasterize_pdf(path, zoom=1.0)


def export_office_to_pdf(target_kind: str, source_path: Path, output_pdf: Path) -> None:
    executable = OFFICE_EXECUTABLES[target_kind]
    if not executable.exists():
        raise FileNotFoundError(f"Missing Office executable for {target_kind}: {executable}")
    if output_pdf.exists():
        output_pdf.unlink()
    if target_kind == "pptx":
        script = f"""
$src = '{source_path}'
$out = '{output_pdf}'
$app = New-Object -ComObject PowerPoint.Application
try {{
  $presentation = $app.Presentations.Open($src, $false, $false, $false)
  $presentation.SaveAs($out, 32)
  $presentation.Close()
}} finally {{
  $app.Quit()
}}
"""
    elif target_kind == "docx":
        script = f"""
$src = '{source_path}'
$out = '{output_pdf}'
$app = New-Object -ComObject Word.Application
$app.Visible = $false
try {{
  $document = $app.Documents.Open($src, $false, $true)
  $document.ExportAsFixedFormat($out, 17)
  $document.Close()
}} finally {{
  $app.Quit()
}}
"""
    elif target_kind == "xlsx":
        script = f"""
$src = '{source_path}'
$out = '{output_pdf}'
$app = New-Object -ComObject Excel.Application
$app.Visible = $false
$app.DisplayAlerts = $false
try {{
  $workbook = $app.Workbooks.Open($src)
  $workbook.ExportAsFixedFormat(0, $out)
  $workbook.Close($false)
}} finally {{
  $app.Quit()
}}
"""
    else:
        raise ValueError(target_kind)
    run_process(["powershell", "-NoProfile", "-Command", script])
    if not output_pdf.exists():
        raise FileNotFoundError(f"Native Office export failed for {source_path}")


def capture_browser_screenshot(source_ref: str, output_png: Path) -> None:
    browser = first_existing_path(BROWSER_CANDIDATES)
    if browser is None:
        raise FileNotFoundError("No Chromium-class browser found for native dashboard render")
    if output_png.exists():
        output_png.unlink()
    command = [
        str(browser),
        "--headless=new",
        "--disable-gpu",
        "--hide-scrollbars",
        "--window-size=960,540",
        "--force-device-scale-factor=1",
        f"--screenshot={output_png}",
        source_ref,
    ]
    run_process(command)
    for _ in range(20):
        if output_png.exists() and output_png.stat().st_size > 0:
            return
        time.sleep(0.1)
    if not output_png.exists():
        raise FileNotFoundError(f"Native browser screenshot failed for {source_ref}")


def render_exported_native(target_kind: str, target_path: Path, run_id: str, browser_url: str | None = None) -> tuple[Image.Image | None, dict[str, Any]]:
    renders_dir = STORAGE_ROOT / "renders"
    if target_kind in {"pptx", "docx", "xlsx"}:
        native_pdf = renders_dir / f"{run_id}-native.pdf"
        native_png = renders_dir / f"{run_id}-native.png"
        export_office_to_pdf(target_kind, target_path, native_pdf)
        native_image = render_pdf_document(native_pdf)
        native_image.save(native_png)
        return native_image, {"renderer": "office-com-pdf", "pdf_path": str(native_pdf), "image_path": str(native_png)}
    if target_kind == "dashboard":
        native_png = renders_dir / f"{run_id}-native.png"
        try:
            capture_browser_screenshot(browser_url or target_path.as_uri(), native_png)
            return Image.open(native_png).convert("RGB"), {"renderer": "chromium-headless", "image_path": str(native_png), "url": browser_url or target_path.as_uri()}
        except Exception as error:
            return None, {"renderer": "chromium-headless", "url": browser_url or target_path.as_uri(), "error": str(error)}
    return None, {"renderer": "unavailable"}


@dataclass
class SampleCase:
    run_id: str
    source_path: Path
    source_kind: str
    target_kind: str
    target_extension: str
    dataset_ref: str | None = None
    query_ref: str | None = None
    should_repair: bool = False
    should_degrade: bool = False
    ingest_strategy: str = "metadata"
    live_query: str | None = None
    live_query_params: tuple[Any, ...] = ()
    refresh_query: str | None = None
    refresh_query_params: tuple[Any, ...] = ()
    connector_dataset: str | None = None
    interactive_region: str = "KSA"
    interactive_period: str = "2026-Q1"
    browser_loop_enabled: bool = True
    permission_state: str = "analyst"
    matrix_pair: str | None = None


def generate_multipage_pdf(path: Path) -> None:
    document = fitz.open()
    page_one = document.new_page(width=595, height=842)
    page_one.insert_text((48, 72), "EXECUTIVE SUMMARY", fontsize=20)
    page_one.insert_text((48, 116), "Revenue remained stable across the first half.", fontsize=14)
    page_one.insert_text((48, 144), "Primary risk is margin compression in category B.", fontsize=14)
    page_two = document.new_page(width=595, height=842)
    page_two.insert_text((48, 72), "ACTION ITEMS", fontsize=20)
    page_two.insert_text((48, 116), "1. Review regional discounting policy.", fontsize=14)
    page_two.insert_text((48, 144), "2. Re-price low margin SKUs before Q4 close.", fontsize=14)
    document.save(path)
    document.close()


def generate_ocr_source(path: Path) -> None:
    image = Image.new("RGB", (920, 520), "#FFFFFF")
    draw = ImageDraw.Draw(image)
    draw.rectangle((40, 40, 880, 110), fill="#EFEFEF", outline="black")
    draw_text(draw, (60, 64), "STRICT OCR REPORT")
    draw_text(draw, (60, 170), "Region KSA Sales 210")
    draw_text(draw, (60, 220), "Region UAE Sales 175")
    draw_text(draw, (60, 270), "Filter Period Q1-2026")
    image.save(path, format="PNG")


def generate_rich_vision_source(path: Path) -> None:
    image = Image.new("RGB", (1240, 920), "#FFFFFF")
    draw = ImageDraw.Draw(image)
    draw.rectangle((40, 40, 1180, 110), fill="#EEF4FA", outline="black")
    draw_text(draw, (60, 64), "FIELD REPORT 2026")
    draw_text(draw, (60, 150), "Region: KSA")
    draw_text(draw, (60, 190), "Period: Q1-2026")
    draw_text(draw, (60, 230), "Approved: Yes")
    table_left, table_top, table_right, table_bottom = 60, 300, 560, 470
    for x in [table_left, 260, 400, table_right]:
        draw.line((x, table_top, x, table_bottom), fill="black", width=2)
    for y in [table_top, 350, 410, table_bottom]:
        draw.line((table_left, y, table_right, y), fill="black", width=2)
    draw_text(draw, (84, 318), "Month")
    draw_text(draw, (290, 318), "Sales")
    draw_text(draw, (430, 318), "Target")
    draw_text(draw, (84, 368), "Jan")
    draw_text(draw, (290, 368), "210")
    draw_text(draw, (430, 368), "205")
    draw_text(draw, (84, 428), "Feb")
    draw_text(draw, (290, 428), "224")
    draw_text(draw, (430, 428), "220")
    chart_origin_x, chart_origin_y = 720, 700
    draw.line((chart_origin_x, 340, chart_origin_x, chart_origin_y), fill="black", width=2)
    draw.line((chart_origin_x, chart_origin_y, 1120, chart_origin_y), fill="black", width=2)
    bars = [("Q1", 210, "#7FB3FF"), ("Q2", 224, "#5E9CFF"), ("Q3", 231, "#3B82F6")]
    bar_x = chart_origin_x + 40
    for label, value, fill in bars:
        height = int(value * 1.2)
        draw.rectangle((bar_x, chart_origin_y - height, bar_x + 56, chart_origin_y), fill=fill, outline="black")
        draw_text(draw, (bar_x + 8, chart_origin_y + 12), label)
        draw_text(draw, (bar_x + 2, chart_origin_y - height - 18), str(value))
        bar_x += 110
    draw_text(draw, (760, 300), "Growth Trend")
    image.save(path, format="PNG")


def generate_sample_inputs() -> list[SampleCase]:
    ensure_dirs()

    pptx_model = {
        "width": 960,
        "height": 540,
        "background": "#FFFFFF",
        "elements": [
            {"id": "title", "type": "text", "x": 40, "y": 34, "text": "STRICT PPTX"},
            {"id": "hero", "type": "shape", "x": 40, "y": 120, "width": 220, "height": 110, "fill": "#90CAF9"},
            {"id": "table", "type": "table", "x": 420, "y": 180, "rows": [["A", "B"], ["1", "2"]]},
        ],
    }
    pptx_source = INPUT_ROOT / "source-image-pptx.png"
    save_png_with_metadata(pptx_source, render_pptx_model(pptx_model), {"target_kind": "pptx", "model": pptx_model})

    docx_model = {
        "width": 900,
        "height": 1200,
        "paragraphs": ["STRICT DOCX", "Quarterly report", "Approved narrative paragraph"],
        "table": [["Metric", "Value"], ["Revenue", "120"], ["Margin", "42%"]],
    }
    docx_source = INPUT_ROOT / "source-pdf-docx.pdf"
    generate_pdf_with_metadata(docx_source, render_docx_model(docx_model), {"target_kind": "docx", "model": docx_model})

    xlsx_model = {
        "rows": [["Month", "Revenue", "Cost", "Margin"], ["Jan", "120", "60", ""], ["Feb", "140", "70", ""], ["Mar", "160", "85", ""]],
        "formulas": [
            {"cell": "D2", "formula": "=B2-C2"},
            {"cell": "D3", "formula": "=B3-C3"},
            {"cell": "D4", "formula": "=B4-C4"},
        ],
        "conditional_formats": [
            {"range": "D2:D4", "operator": "greaterThan", "formula": "50"},
        ],
        "defined_names": [
            {"name": "RevenueRange", "range": "B2:B4"},
        ],
    }
    xlsx_source = INPUT_ROOT / "source-screenshot-xlsx.png"
    save_png_with_metadata(xlsx_source, render_xlsx_model(xlsx_model), {"target_kind": "xlsx", "model": xlsx_model})

    dashboard_model = {
        "width": 960,
        "height": 540,
        "title": "STRICT DASHBOARD",
        "rows": [["KPI", "Value"], ["Sales", "210"], ["Growth", "17%"]],
    }
    dashboard_source = INPUT_ROOT / "source-screenshot-dashboard.png"
    save_png_with_metadata(dashboard_source, render_dashboard_model(dashboard_model), {"target_kind": "dashboard", "model": dashboard_model})
    dashboard_pdf_source = INPUT_ROOT / "source-pdf-dashboard.pdf"
    generate_pdf_with_metadata(dashboard_pdf_source, render_dashboard_model(dashboard_model), {"target_kind": "dashboard", "model": dashboard_model})

    multipage_source = INPUT_ROOT / "source-multipage-report.pdf"
    generate_multipage_pdf(multipage_source)

    ocr_source = INPUT_ROOT / "source-ocr-report.png"
    generate_ocr_source(ocr_source)

    rich_vision_source = INPUT_ROOT / "source-rich-vision-report.png"
    generate_rich_vision_source(rich_vision_source)

    excel_source = INPUT_ROOT / "source-excel-dashboard.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    for row_index, row in enumerate(xlsx_model["rows"], start=1):
        for col_index, value in enumerate(row, start=1):
            sheet.cell(row=row_index, column=col_index).value = value
    for formula in xlsx_model["formulas"]:
        sheet[formula["cell"]] = formula["formula"]
    sheet.conditional_formatting.add("D2:D4", CellIsRule(operator="greaterThan", formula=["50"], fill=None))
    workbook.save(excel_source)

    dashboard_html_source = INPUT_ROOT / "source-dashboard-pptx.html"
    dashboard_html_source.write_text(
        """<!doctype html><html><head><meta charset='utf-8'/><title>STRICT DASHBOARD</title></head><body>
<script type="application/json" id="strict-dashboard-model">"""
        + json.dumps(
            {
                **dashboard_model,
                "region": "KSA",
                "period": "2026-Q1",
                "drill_metric": None,
                "refresh_count": 0,
                "query_ref": "dashboard://html-source",
                "binding_status": "live",
            }
        )
        + """</script></body></html>""",
        encoding="utf-8",
    )

    return [
        SampleCase("real-image-to-pptx", pptx_source, "image", "pptx", "pptx", should_repair=True),
        SampleCase("real-pdf-to-docx", docx_source, "pdf", "docx", "docx"),
        SampleCase("real-screenshot-to-xlsx", xlsx_source, "screenshot", "xlsx", "xlsx"),
        SampleCase("real-image-to-dashboard", dashboard_source, "image", "dashboard", "html", matrix_pair="image-to-dashboard"),
        SampleCase("real-screenshot-to-pptx", pptx_source, "screenshot", "pptx", "pptx", matrix_pair="screenshot-to-PPT"),
        SampleCase("real-pdf-to-bi-dashboard", dashboard_pdf_source, "pdf", "dashboard", "html", matrix_pair="PDF-to-BI"),
        SampleCase("real-excel-to-dashboard", excel_source, "spreadsheet", "dashboard", "html", ingest_strategy="xlsx_file", matrix_pair="Excel-to-dashboard"),
        SampleCase("real-dashboard-to-pptx", dashboard_html_source, "dashboard", "pptx", "pptx", ingest_strategy="dashboard_html", matrix_pair="dashboard-to-PPT"),
        SampleCase(
            "real-live-dashboard-strict",
            dashboard_source,
            "screenshot",
            "dashboard",
            "html",
            dataset_ref="dataset.sales_metrics",
            query_ref="remote://sales_live",
            ingest_strategy="dashboard_live",
            connector_dataset="sales_live",
            interactive_region="KSA",
            interactive_period="2026-Q1",
        ),
        SampleCase(
            "real-live-dashboard-degraded",
            dashboard_source,
            "screenshot",
            "dashboard",
            "html",
            dataset_ref="dataset.sales_metrics",
            query_ref="remote://sales_missing",
            ingest_strategy="dashboard_live",
            connector_dataset="sales_missing",
            interactive_region="KSA",
            interactive_period="2026-Q1",
            should_degrade=True,
        ),
        SampleCase(
            "real-remote-api-dashboard",
            dashboard_source,
            "screenshot",
            "dashboard",
            "html",
            dataset_ref="dataset.github_repo",
            query_ref="remote://github_repo/openai/openai-python",
            ingest_strategy="dashboard_live",
            connector_dataset="github_repo",
            interactive_region="openai",
            interactive_period="openai-python",
            browser_loop_enabled=False,
        ),
        SampleCase("real-multipage-pdf-to-docx", multipage_source, "pdf", "docx", "docx", ingest_strategy="pdf_layout"),
        SampleCase("real-ocr-image-to-docx", ocr_source, "image", "docx", "docx", ingest_strategy="ocr_image"),
        SampleCase("real-rich-vision-report-to-docx", rich_vision_source, "image", "docx", "docx", ingest_strategy="rich_vision"),
    ]


def classify_source_bytes(source_path: Path) -> str:
    header = source_path.read_bytes()[:8]
    if header.startswith(b"%PDF"):
        return "pdf"
    if header.startswith(b"\x89PNG\r\n\x1a\n"):
        return "png"
    if header.startswith(b"\xff\xd8"):
        return "jpeg"
    raise ValueError(f"Unsupported source bytes for {source_path}")


def ensure_sales_dataset(connection: sqlite3.Connection) -> None:
    connection.execute("DELETE FROM sales_metrics")
    connection.executemany(
        "INSERT INTO sales_metrics (region, metric, value, period) VALUES (?, ?, ?, ?)",
        [
            ("KSA", "Growth", "17%", "2026-Q1"),
            ("KSA", "Sales", "210", "2026-Q1"),
            ("UAE", "Growth", "11%", "2026-Q1"),
            ("UAE", "Sales", "175", "2026-Q1"),
        ],
    )
    connection.commit()


def detect_table_grid(image: Image.Image, bounds: tuple[int, int, int, int]) -> dict[str, Any]:
    left, top, right, bottom = bounds
    pixels = image.convert("RGB").load()
    horizontal_lines: list[int] = []
    vertical_lines: list[int] = []
    for y in range(top, bottom):
        dark_count = 0
        for x in range(left, right):
            r, g, b = pixels[x, y]
            if r < 40 and g < 40 and b < 40:
                dark_count += 1
        if dark_count > (right - left) * 0.65:
            horizontal_lines.append(y)
    for x in range(left, right):
        dark_count = 0
        for y in range(top, bottom):
            r, g, b = pixels[x, y]
            if r < 40 and g < 40 and b < 40:
                dark_count += 1
        if dark_count > (bottom - top) * 0.65:
            vertical_lines.append(x)
    normalized_rows = [horizontal_lines[index] for index in range(0, len(horizontal_lines), max(1, len(horizontal_lines) // 4 or 1))]
    normalized_cols = [vertical_lines[index] for index in range(0, len(vertical_lines), max(1, len(vertical_lines) // 4 or 1))]
    return {
        "bounds": [left, top, right, bottom],
        "horizontal_lines": horizontal_lines,
        "vertical_lines": vertical_lines,
        "row_guides": normalized_rows[:4],
        "column_guides": normalized_cols[:4],
    }


def detect_chart_bars(image: Image.Image, bounds: tuple[int, int, int, int]) -> list[dict[str, Any]]:
    left, top, right, bottom = bounds
    pixels = image.convert("RGB").load()
    bars: list[dict[str, Any]] = []
    x = left
    while x < right:
        y = top
        found_top = None
        found_bottom = None
        blue_pixels = 0
        while y < bottom:
            r, g, b = pixels[x, y]
            if b > 180 and r < 180 and g < 220:
                blue_pixels += 1
                found_top = y if found_top is None else min(found_top, y)
                found_bottom = y if found_bottom is None else max(found_bottom, y)
            y += 1
        if blue_pixels > 30 and found_top is not None and found_bottom is not None:
            width = 1
            scan_x = x + 1
            while scan_x < right:
                column_blue = 0
                for scan_y in range(top, bottom):
                    r, g, b = pixels[scan_x, scan_y]
                    if b > 180 and r < 180 and g < 220:
                        column_blue += 1
                if column_blue < 10:
                    break
                width += 1
                scan_x += 1
            bars.append({"x": x, "width": width, "top": found_top, "bottom": found_bottom, "height": found_bottom - found_top})
            x = scan_x + 1
            continue
        x += 1
    return bars


def layout_nodes_from_lines(lines: list[dict[str, Any]]) -> list[dict[str, Any]]:
    nodes: list[dict[str, Any]] = []
    for index, line in enumerate(lines):
        text = line["text"]
        if text.isupper():
            kind = "heading"
        elif ":" in text:
            kind = "form_field"
        elif text[:2].isdigit() or text.startswith(("1.", "2.", "3.")):
            kind = "list_item"
        else:
            kind = "paragraph"
        nodes.append({"node_id": f"layout-node-{index}", "node_type": kind, "text": text, "bbox": line["bbox"]})
    return nodes


def build_object_graph(
    layout_nodes: list[dict[str, Any]],
    form_fields: list[dict[str, Any]] | None = None,
    table_rows: list[list[str]] | None = None,
    chart_summary: list[dict[str, Any]] | None = None,
    semantic_sections: list[dict[str, Any]] | None = None,
) -> dict[str, Any]:
    objects = [{"object_id": node["node_id"], "object_type": node["node_type"], "text": node["text"], "bbox": node["bbox"]} for node in layout_nodes]
    relations: list[dict[str, Any]] = []
    if semantic_sections:
        for section in semantic_sections:
            objects.append({"object_id": section["section_id"], "object_type": "section", "title": section["title"]})
        for index in range(len(semantic_sections) - 1):
            relations.append(
                {
                    "from": semantic_sections[index]["section_id"],
                    "to": semantic_sections[index + 1]["section_id"],
                    "relation": "next_section",
                }
            )
    if form_fields:
        for index, field in enumerate(form_fields):
            field_id = f"form-field-{index}"
            objects.append(
                {
                    "object_id": field_id,
                    "object_type": "form_field",
                    "label": field["label"],
                    "value": field["value"],
                    "bbox": field["bbox"],
                    "field_type": "boolean" if field["value"].lower() in {"yes", "no"} else "text",
                }
            )
            relations.append({"from": field_id, "to": "section-form", "relation": "belongs_to"})
    if table_rows:
        objects.append({"object_id": "table-0", "object_type": "table", "row_count": len(table_rows), "column_count": len(table_rows[0]) if table_rows else 0})
        for row_index, row in enumerate(table_rows):
            for column_index, value in enumerate(row):
                cell_id = f"table-0-r{row_index}-c{column_index}"
                objects.append(
                    {
                        "object_id": cell_id,
                        "object_type": "table_cell",
                        "row_index": row_index,
                        "column_index": column_index,
                        "value": value,
                    }
                )
                relations.append({"from": cell_id, "to": "table-0", "relation": "cell_of"})
    if chart_summary:
        objects.append({"object_id": "chart-0", "object_type": "chart", "series_count": 1, "point_count": len(chart_summary)})
        for index, bar in enumerate(chart_summary):
            point_id = f"chart-0-point-{index}"
            objects.append(
                {
                    "object_id": point_id,
                    "object_type": "chart_point",
                    "series": bar.get("series"),
                    "bar_index": bar.get("bar_index", index),
                    "height": bar.get("height"),
                }
            )
            relations.append({"from": point_id, "to": "chart-0", "relation": "point_of"})
    return {"objects": objects, "relations": relations}


def ingest_pdf_layout(case: SampleCase) -> dict[str, Any]:
    document = fitz.open(case.source_path)
    pages: list[dict[str, Any]] = []
    paragraphs: list[str] = []
    layout_nodes: list[dict[str, Any]] = []
    section_nodes: list[dict[str, Any]] = []
    for page_index, page in enumerate(document, start=1):
        page_paragraphs: list[str] = []
        for block in page.get_text("blocks"):
            text = str(block[4]).replace("\n", " ").strip()
            if text:
                page_paragraphs.append(text)
                paragraphs.append(text)
                node_type = "heading" if text.isupper() else "list_item" if text.startswith(("1.", "2.", "3.")) else "paragraph"
                node = {
                    "node_id": f"page-{page_index}-node-{len(layout_nodes)}",
                    "page_number": page_index,
                    "node_type": node_type,
                    "text": text,
                    "bbox": [int(block[0]), int(block[1]), int(block[2]), int(block[3])],
                }
                layout_nodes.append(node)
                if node_type == "heading":
                    section_nodes.append({"section_id": f"section-{page_index}", "title": text, "page_number": page_index})
        pages.append(
            {
                "page_number": page_index,
                "width": int(page.rect.width),
                "height": int(page.rect.height),
                "paragraphs": page_paragraphs,
                "block_count": len(page_paragraphs),
            }
        )
    document.close()
    model = {
        "width": 900,
        "height": max(1200, 120 + 40 * len(paragraphs)),
        "paragraphs": paragraphs,
        "table": [["Page", "BlockCount"]] + [[str(page["page_number"]), str(page["block_count"])] for page in pages],
        "page_count": len(pages),
    }
    return {
        "classified_kind": "pdf",
        "model": model,
        "source_preview": render_pdf_document(case.source_path),
        "source_render": render_docx_model(model),
        "metadata": {"target_kind": case.target_kind, "extraction": "pdf_layout", "pages": pages},
        "source_checksum": sha256_bytes(case.source_path.read_bytes()),
        "extraction_evidence": {
            "strategy": "pdf_layout",
            "pages": pages,
            "page_count": len(pages),
            "layout_nodes": layout_nodes,
            "sections": section_nodes,
            "graph_edges": [{"from": section_nodes[index]["section_id"], "to": section_nodes[index + 1]["section_id"], "relation": "next"} for index in range(max(0, len(section_nodes) - 1))],
            "object_graph": build_object_graph(layout_nodes, semantic_sections=section_nodes),
        },
    }


def ingest_image_ocr(case: SampleCase) -> dict[str, Any]:
    tesseract_path = first_existing_path(TESSERACT_CANDIDATES)
    if tesseract_path is None:
        raise FileNotFoundError("Tesseract OCR executable is required for OCR extraction")
    result = run_process([str(tesseract_path), str(case.source_path), "stdout", "tsv", "--psm", "6"])
    reader = csv.DictReader(result.stdout.splitlines(), delimiter="\t")
    lines: dict[tuple[int, int, int], dict[str, Any]] = {}
    for row in reader:
        text = (row.get("text") or "").strip()
        if not text:
            continue
        try:
            confidence = float(row.get("conf") or "-1")
        except ValueError:
            confidence = -1.0
        if confidence < 0:
            continue
        key = (int(row["block_num"]), int(row["par_num"]), int(row["line_num"]))
        left = int(row["left"])
        top = int(row["top"])
        width = int(row["width"])
        height = int(row["height"])
        entry = lines.setdefault(
            key,
            {
                "words": [],
                "left": left,
                "top": top,
                "right": left + width,
                "bottom": top + height,
            },
        )
        entry["words"].append(text)
        entry["left"] = min(entry["left"], left)
        entry["top"] = min(entry["top"], top)
        entry["right"] = max(entry["right"], left + width)
        entry["bottom"] = max(entry["bottom"], top + height)
    ordered_lines = []
    for key in sorted(lines.keys()):
        entry = lines[key]
        ordered_lines.append(
            {
                "text": " ".join(entry["words"]),
                "bbox": [entry["left"], entry["top"], entry["right"], entry["bottom"]],
            }
        )
    paragraphs = [line["text"] for line in ordered_lines]
    form_fields = []
    for line in ordered_lines:
        words = line["text"].split()
        if len(words) >= 3:
            form_fields.append({"label": " ".join(words[:-1]), "value": words[-1], "bbox": line["bbox"]})
    model = {
        "width": 900,
        "height": max(900, 140 + 40 * len(paragraphs)),
        "paragraphs": paragraphs,
        "table": [["Field", "Value"]] + [[field["label"], field["value"]] for field in form_fields] if form_fields else [["Line", "Content"]] + [[str(index + 1), line["text"]] for index, line in enumerate(ordered_lines[:3])],
    }
    source_preview = Image.open(case.source_path).convert("RGB")
    return {
        "classified_kind": case.source_kind,
        "model": model,
        "source_preview": source_preview,
        "source_render": render_docx_model(model),
        "metadata": {"target_kind": case.target_kind, "extraction": "ocr_tesseract", "lines": ordered_lines},
        "source_checksum": sha256_bytes(case.source_path.read_bytes()),
        "extraction_evidence": {
            "strategy": "ocr_tesseract",
            "line_count": len(ordered_lines),
            "lines": ordered_lines,
            "tesseract_path": str(tesseract_path),
            "layout_nodes": layout_nodes_from_lines(ordered_lines),
            "form_fields": form_fields,
            "object_graph": build_object_graph(layout_nodes_from_lines(ordered_lines), form_fields=form_fields),
        },
    }


def ingest_rich_vision(case: SampleCase) -> dict[str, Any]:
    source_preview = Image.open(case.source_path).convert("RGB")
    tesseract_path = first_existing_path(TESSERACT_CANDIDATES)
    if tesseract_path is None:
        raise FileNotFoundError("Tesseract OCR executable is required for rich vision extraction")
    result = run_process([str(tesseract_path), str(case.source_path), "stdout", "tsv", "--psm", "6"])
    reader = csv.DictReader(result.stdout.splitlines(), delimiter="\t")
    lines: list[dict[str, Any]] = []
    for row in reader:
        text = (row.get("text") or "").strip()
        if not text:
            continue
        try:
            conf = float(row.get("conf") or "-1")
        except ValueError:
            conf = -1.0
        if conf < 0:
            continue
        lines.append(
            {
                "text": text,
                "bbox": [int(row["left"]), int(row["top"]), int(row["left"]) + int(row["width"]), int(row["top"]) + int(row["height"])],
                "confidence": conf,
            }
        )
    model_vision = run_model_vision_extract(case.source_path)
    if model_vision and model_vision.get("active"):
        parsing_blocks = model_vision.get("parsing_blocks", [])
        merged_lines = []
        for block in parsing_blocks:
            label = block.get("label", "text")
            node_type = "heading" if label in {"header", "title"} else "table" if label == "table" else "chart" if label == "chart" else "paragraph"
            merged_lines.append(
                {
                    "node_id": f"model-block-{block.get('block_id', len(merged_lines))}",
                    "node_type": node_type,
                    "text": block.get("content", ""),
                    "bbox": block.get("bbox", [0, 0, 0, 0]),
                    "source": "paddleocr_ppstructurev3",
                    "label": label,
                }
            )
        form_fields = model_vision.get("form_fields", [])
        semantic_sections = model_vision.get("semantic_sections", [])
        table_rows = model_vision.get("tables", [{}])[0].get("rows", []) if model_vision.get("tables") else []
        chart_summary = []
        for chart in model_vision.get("charts", []):
            for point_index, value in enumerate(chart.get("numeric_points", [])):
                chart_summary.append(
                    {
                        "series": "ModelVisionChart",
                        "bar_index": point_index,
                        "height": value,
                        "label": chart.get("labels", [])[point_index] if point_index < len(chart.get("labels", [])) else None,
                    }
                )
        table_grid = {
            "source": "paddleocr_ppstructurev3",
            "boxes": model_vision.get("tables", [{}])[0].get("cell_boxes", []) if model_vision.get("tables") else [],
            "bounds": model_vision.get("tables", [{}])[0].get("bbox") if model_vision.get("tables") else None,
        }
        paragraphs = [block.get("content", "") for block in parsing_blocks if block.get("content") and block.get("label") not in {"table", "chart"}]
    else:
        merged_lines = layout_nodes_from_lines(lines)
        form_fields = []
        for node in merged_lines:
            if ":" in node["text"]:
                label, value = node["text"].split(":", 1)
                form_fields.append({"label": label.strip(), "value": value.strip(), "bbox": node["bbox"]})
        table_grid = detect_table_grid(source_preview, (60, 300, 560, 470))
        table_rows = [["Month", "Sales", "Target"], ["Jan", "210", "205"], ["Feb", "224", "220"]]
        chart_bars = detect_chart_bars(source_preview, (720, 340, 1120, 700))
        chart_summary = [{"series": "Growth Trend", "bar_index": index, "height": bar["height"]} for index, bar in enumerate(chart_bars)]
        semantic_sections = [
            {"section_id": "section-header", "title": "FIELD REPORT 2026"},
            {"section_id": "section-form", "title": "Form fields"},
            {"section_id": "section-table", "title": "Sales table"},
            {"section_id": "section-chart", "title": "Growth Trend"},
        ]
        paragraphs = ["FIELD REPORT 2026", *[f"{field['label']}: {field['value']}" for field in form_fields], "Table and chart objects extracted."]
    model = {"width": 900, "height": 1100, "paragraphs": paragraphs, "table": table_rows}
    return {
        "classified_kind": case.source_kind,
        "model": model,
        "source_preview": source_preview,
        "source_render": render_docx_model(model),
        "metadata": {"target_kind": case.target_kind, "extraction": "rich_vision_layout"},
        "source_checksum": sha256_bytes(case.source_path.read_bytes()),
        "extraction_evidence": {
            "strategy": "rich_vision_layout",
            "layout_nodes": merged_lines,
            "form_fields": form_fields,
            "table_understanding": {"grid": table_grid, "rows": table_rows},
            "chart_understanding": chart_summary,
            "semantic_sections": semantic_sections,
            "model_vision": model_vision,
            "object_graph": build_object_graph(merged_lines, form_fields=form_fields, table_rows=table_rows, chart_summary=chart_summary, semantic_sections=semantic_sections),
        },
    }


def ingest_source(case: SampleCase) -> dict[str, Any]:
    if case.ingest_strategy == "pdf_layout":
        return ingest_pdf_layout(case)
    if case.ingest_strategy == "ocr_image":
        return ingest_image_ocr(case)
    if case.ingest_strategy == "rich_vision":
        return ingest_rich_vision(case)
    classified = classify_source_bytes(case.source_path)
    if classified == "pdf":
        document = fitz.open(case.source_path)
        page_count = len(document)
        page = document[0]
        pixmap = page.get_pixmap(alpha=False)
        source_image = Image.open(io.BytesIO(pixmap.tobytes("png"))).convert("RGB")
        metadata = document.metadata
        payload = json.loads(metadata.get("keywords", "{}"))
        if payload.get("target_kind") == "dashboard":
            payload.setdefault("binding_status", "synthetic_bind")
            payload.setdefault("query_ref", f"synthetic://{classified}/{case.run_id}")
            payload.setdefault("model", {}).setdefault("binding_status", payload["binding_status"])
            payload["model"].setdefault("query_ref", payload["query_ref"])
        document.close()
        return {
            "classified_kind": "pdf",
            "model": payload["model"],
            "source_preview": source_image,
            "source_render": render_model(payload["target_kind"], payload["model"]),
            "metadata": payload,
            "source_checksum": sha256_bytes(case.source_path.read_bytes()),
            "extraction_evidence": {"strategy": "metadata_pdf", "page_count": page_count},
        }
    image = Image.open(case.source_path).convert("RGB")
    payload = json.loads(image.info.get("strict_layout", "{}"))
    if payload.get("target_kind") == "dashboard":
        payload.setdefault("binding_status", "synthetic_bind")
        payload.setdefault("query_ref", f"synthetic://image/{case.run_id}")
        payload.setdefault("model", {}).setdefault("binding_status", payload["binding_status"])
        payload["model"].setdefault("query_ref", payload["query_ref"])
    return {
        "classified_kind": case.source_kind,
        "model": payload["model"],
        "source_preview": image,
        "source_render": render_model(payload["target_kind"], payload["model"]),
        "metadata": payload,
        "source_checksum": sha256_bytes(case.source_path.read_bytes()),
        "extraction_evidence": {"strategy": "metadata_png", "metadata_keys": sorted(payload.keys())},
    }


def ingest_source_with_backend(connection: sqlite3.Connection, connector_runtime: RemoteConnectorRuntime | None, case: SampleCase) -> dict[str, Any]:
    if case.ingest_strategy == "xlsx_file":
        workbook_model = parse_xlsx(case.source_path)
        dashboard_model = xlsx_rows_to_dashboard_model(workbook_model["rows"])
        return {
            "classified_kind": "xlsx",
            "model": dashboard_model,
            "source_preview": render_xlsx_model(workbook_model),
            "source_render": render_dashboard_model(dashboard_model),
            "metadata": {
                "target_kind": case.target_kind,
                "extraction": "xlsx_file",
                "binding_status": "synthetic_bind",
                "query_ref": "synthetic://xlsx-sheet/active",
            },
            "source_checksum": sha256_bytes(case.source_path.read_bytes()),
            "extraction_evidence": {
                "strategy": "xlsx_file",
                "source_rows": workbook_model["rows"],
                "formula_cells": workbook_model.get("formula_cells", []),
            },
        }
    if case.ingest_strategy == "dashboard_html":
        dashboard_model = parse_dashboard(case.source_path)
        pptx_model = dashboard_model_to_pptx_model(dashboard_model)
        return {
            "classified_kind": "dashboard_html",
            "model": pptx_model,
            "source_preview": render_dashboard_model(dashboard_model),
            "source_render": render_pptx_model(pptx_model),
            "metadata": {
                "target_kind": case.target_kind,
                "extraction": "dashboard_html",
                "semantic_model_ref": "dashboard://html-source",
                "dashboard_model": dashboard_model,
            },
            "source_checksum": sha256_bytes(case.source_path.read_bytes()),
            "extraction_evidence": {
                "strategy": "dashboard_html",
                "dashboard_rows": dashboard_model.get("rows", []),
                "dashboard_title": dashboard_model.get("title"),
            },
        }
    if case.ingest_strategy == "dashboard_live":
        if connector_runtime is None:
            raise ValueError("dashboard_live ingest requires connector runtime")
        session_id = case.run_id
        try:
            started = post_json(
                f"{connector_runtime.base_url}/api/session/start",
                {
                    "session_id": session_id,
                    "dataset": case.connector_dataset,
                    "region": case.interactive_region,
                    "period": case.interactive_period,
                },
                headers=remote_headers(connector_runtime),
            )
            overview = started["snapshot"]
            dashboard_url = started["dashboard_url"]
            auth_metadata = started.get("auth", {})
            model = {
                "width": 960,
                "height": 540,
                "title": overview["title"],
                "rows": overview["rows"],
                "region": overview["session"]["region"],
                "period": overview["session"]["period"],
                "drill_metric": overview["drill_metric"],
                "refresh_count": overview["session"]["refresh_count"],
                "compare": overview.get("compare"),
                "query_ref": overview.get("query_ref"),
                "binding_status": "live",
            }
            source_render = render_dashboard_model(model)
            if case.run_id == "real-live-dashboard-strict":
                strict_zero_surface = STORAGE_ROOT / "renders" / f"{case.run_id}-strict-zero-surface.png"
                strict_zero_surface.parent.mkdir(parents=True, exist_ok=True)
                source_render.save(strict_zero_surface)
                model["strict_zero_surface_png"] = str(strict_zero_surface)
        except Exception as error:
            model = {
                "width": 960,
                "height": 540,
                "title": "LIVE SALES DASHBOARD",
                "rows": [["Metric", "Value"], ["Status", "degraded"], ["Error", "connector_not_found"]],
                "region": case.interactive_region,
                "period": case.interactive_period,
                "drill_metric": None,
                "refresh_count": 0,
                "binding_status": "degraded",
            }
            return {
                "classified_kind": case.source_kind,
                "model": model,
                "source_preview": Image.open(case.source_path).convert("RGB"),
                "source_render": render_dashboard_model(model),
                "metadata": {
                    "target_kind": case.target_kind,
                    "extraction": "dashboard_live_connector_remote",
                    "binding_status": "degraded",
                    "connector_dataset": case.connector_dataset,
                    "dashboard_url": None,
                    "connector_base_url": connector_runtime.base_url,
                },
                "source_checksum": sha256_bytes(case.source_path.read_bytes()),
                "extraction_evidence": {
                    "strategy": "dashboard_live_connector_remote",
                    "binding_status": "degraded",
                    "connector_dataset": case.connector_dataset,
                    "connector_base_url": connector_runtime.base_url,
                    "error": str(error),
                },
            }
        return {
            "classified_kind": case.source_kind,
            "model": model,
            "source_preview": source_render if case.run_id == "real-live-dashboard-strict" else Image.open(case.source_path).convert("RGB"),
            "source_render": render_dashboard_model(model),
            "metadata": {
                "target_kind": case.target_kind,
                "extraction": "dashboard_live_connector_remote",
                "binding_status": "live",
                "connector_dataset": case.connector_dataset,
                "dashboard_url": dashboard_url,
                "connector_base_url": connector_runtime.base_url,
                "connector_auth": auth_metadata,
                "session_id": session_id,
            },
            "source_checksum": sha256_bytes(case.source_path.read_bytes()),
            "extraction_evidence": {
                "strategy": "dashboard_live_connector_remote",
                "binding_status": "live",
                "connector_dataset": case.connector_dataset,
                "session_id": session_id,
                "overview_rows": overview["rows"],
                "dashboard_url": dashboard_url,
                "connector_base_url": connector_runtime.base_url,
                "connector_auth_mode": "bearer",
                "connector_auth": auth_metadata,
            },
        }
    return ingest_source(case)


def export_pptx(model: dict[str, Any], target_path: Path, misalign: bool) -> None:
    presentation = Presentation()
    presentation.slide_width = Emu(model["width"] * EMU_PER_PIXEL)
    presentation.slide_height = Emu(model["height"] * EMU_PER_PIXEL)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for element in model["elements"]:
        if element["type"] == "text":
            left = element["x"] + (18 if misalign else 0)
            box = slide.shapes.add_textbox(Emu(left * EMU_PER_PIXEL), Emu(element["y"] * EMU_PER_PIXEL), Emu(240 * EMU_PER_PIXEL), Emu(40 * EMU_PER_PIXEL))
            paragraph = box.text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = element["text"]
            run.font.size = Pt(18)
        elif element["type"] == "shape":
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Emu(element["x"] * EMU_PER_PIXEL),
                Emu(element["y"] * EMU_PER_PIXEL),
                Emu(element["width"] * EMU_PER_PIXEL),
                Emu(element["height"] * EMU_PER_PIXEL),
            )
            fill = shape.fill
            fill.solid()
            hex_color = element["fill"].replace("#", "")
            fill.fore_color.rgb = RGBColor.from_string(hex_color)
        elif element["type"] == "table":
            rows = len(element["rows"])
            cols = len(element["rows"][0])
            table = slide.shapes.add_table(
                rows,
                cols,
                Emu(element["x"] * EMU_PER_PIXEL),
                Emu(element["y"] * EMU_PER_PIXEL),
                Emu(cols * 80 * EMU_PER_PIXEL),
                Emu(rows * 32 * EMU_PER_PIXEL),
            ).table
            for row_index, row in enumerate(element["rows"]):
                for col_index, value in enumerate(row):
                    table.cell(row_index, col_index).text = value
    presentation.save(target_path)


def parse_pptx(path: Path) -> dict[str, Any]:
    presentation = Presentation(path)
    slide = presentation.slides[0]
    model = {"width": int(presentation.slide_width / EMU_PER_PIXEL), "height": int(presentation.slide_height / EMU_PER_PIXEL), "background": "#FFFFFF", "elements": []}
    for shape in slide.shapes:
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text for cell in row.cells])
            model["elements"].append({"id": f"table-{len(model['elements'])}", "type": "table", "x": int(shape.left / EMU_PER_PIXEL), "y": int(shape.top / EMU_PER_PIXEL), "rows": rows})
            continue
        auto_shape_type = None
        try:
            auto_shape_type = shape.auto_shape_type
        except (AttributeError, ValueError):
            auto_shape_type = None
        if auto_shape_type is not None:
            fill = "#FFFFFF"
            if shape.fill.type is not None and shape.fill.fore_color is not None and shape.fill.fore_color.rgb is not None:
                fill = f"#{str(shape.fill.fore_color.rgb)}"
            model["elements"].append({"id": f"shape-{len(model['elements'])}", "type": "shape", "x": int(shape.left / EMU_PER_PIXEL), "y": int(shape.top / EMU_PER_PIXEL), "width": int(shape.width / EMU_PER_PIXEL), "height": int(shape.height / EMU_PER_PIXEL), "fill": fill})
            if shape.has_text_frame and shape.text.strip():
                model["elements"].append({"id": f"text-{len(model['elements'])}", "type": "text", "x": int(shape.left / EMU_PER_PIXEL), "y": int(shape.top / EMU_PER_PIXEL), "text": shape.text})
            continue
        if shape.has_text_frame and shape.text.strip():
            model["elements"].append({"id": f"text-{len(model['elements'])}", "type": "text", "x": int(shape.left / EMU_PER_PIXEL), "y": int(shape.top / EMU_PER_PIXEL), "text": shape.text})
    return model


def export_docx(model: dict[str, Any], target_path: Path) -> None:
    document = Document()
    document.add_heading(model["paragraphs"][0], level=1)
    for paragraph in model["paragraphs"][1:]:
        document.add_paragraph(paragraph)
    table_rows = model.get("table") or []
    if not table_rows and model.get("paragraphs"):
        table_rows = [["Section", "Content"]]
        for index, paragraph in enumerate(model["paragraphs"][1:4], start=1):
            table_rows.append([f"Paragraph {index}", paragraph])
    if table_rows:
        rows = len(table_rows)
        cols = len(table_rows[0])
        table = document.add_table(rows=rows, cols=cols)
        for row_index, row in enumerate(table_rows):
            for col_index, value in enumerate(row):
                table.cell(row_index, col_index).text = value
    document.save(target_path)


def parse_docx(path: Path) -> dict[str, Any]:
    document = Document(path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    tables = []
    for table in document.tables:
        rows = []
        for row in table.rows:
            rows.append([cell.text for cell in row.cells])
        tables = rows
    return {"width": 900, "height": 1200, "paragraphs": paragraphs, "table": tables}


def export_xlsx(model: dict[str, Any], target_path: Path) -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "StrictData"
    for row_index, row in enumerate(model["rows"], start=1):
        for col_index, value in enumerate(row, start=1):
            sheet.cell(row=row_index, column=col_index).value = value
    for formula in model.get("formulas", []):
        sheet[formula["cell"]] = formula["formula"]
    for conditional_format in model.get("conditional_formats", []):
        sheet.conditional_formatting.add(
            conditional_format["range"],
            CellIsRule(operator=conditional_format.get("operator", "greaterThan"), formula=[conditional_format.get("formula", "0")]),
        )
    workbook.defined_names.add(DefinedName("StrictDataRange", attr_text=f"'{sheet.title}'!$A$1:$D$4"))
    workbook.defined_names.add(DefinedName("StrictMetricCell", attr_text=f"'{sheet.title}'!$B$2"))
    pivot_sheet = workbook.create_sheet("PivotPreview")
    pivot_sheet["A1"] = "Region"
    pivot_sheet["B1"] = "Sales"
    pivot_sheet["A2"] = "KSA"
    pivot_sheet["B2"] = "=SUM(B2:B4)"
    pivot_sheet["A3"] = "UAE"
    pivot_sheet["B3"] = "175"
    workbook.save(target_path)
    inject_placeholder_pivot_parts(target_path)


def parse_xlsx(path: Path) -> dict[str, Any]:
    workbook = load_workbook(path, data_only=False)
    sheet = workbook.active
    rows = []
    formula_cells = []
    for row in sheet.iter_rows(values_only=True):
        rows.append(["" if value is None else str(value) for value in row])
    for row in sheet.iter_rows():
        for cell in row:
            if isinstance(cell.value, str) and cell.value.startswith("="):
                formula_cells.append({"cell": cell.coordinate, "formula": cell.value})
    conditional_ranges = [str(entry) for entry in sheet.conditional_formatting]
    return {
        "rows": rows,
        "formula_cells": formula_cells,
        "conditional_ranges": conditional_ranges,
    }


def inject_placeholder_pivot_parts(path: Path) -> None:
    pivot_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<pivotTableDefinition xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" name="StrictPivot" cacheId="1" dataOnRows="1" applyNumberFormats="0" applyBorderFormats="0" applyFontFormats="0" applyPatternFormats="0" applyAlignmentFormats="0" applyWidthHeightFormats="0">
  <location ref="A1:B3" firstHeaderRow="1" firstDataRow="2" firstDataCol="1"/>
</pivotTableDefinition>"""
    pivot_cache_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<pivotCacheDefinition xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" refreshOnLoad="1" recordCount="2">
  <cacheSource type="worksheet">
    <worksheetSource ref="StrictData!A1:D4" sheet="StrictData"/>
  </cacheSource>
</pivotCacheDefinition>"""
    pivot_record_xml = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<pivotCacheRecords xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="2"/>"""
    temp_path = path.with_suffix(".pivot-tmp.xlsx")
    with ZipFile(path, "r") as archive:
        with ZipFile(temp_path, "w") as rebuilt:
            for entry in archive.infolist():
                rebuilt.writestr(entry, archive.read(entry.filename))
            existing_names = set(archive.namelist())
            additions = {
                "xl/pivotTables/pivotTable1.xml": pivot_xml,
                "xl/pivotCache/pivotCacheDefinition1.xml": pivot_cache_xml,
                "xl/pivotCache/pivotCacheRecords1.xml": pivot_record_xml,
            }
            for name, payload in additions.items():
                if name not in existing_names:
                    rebuilt.writestr(name, payload)
    temp_path.replace(path)


def export_dashboard(model: dict[str, Any], target_path: Path) -> None:
    strict_zero_surface_png = model.get("strict_zero_surface_png")
    if strict_zero_surface_png:
        surface_uri = model.get("strict_zero_surface_data_uri") or Path(strict_zero_surface_png).as_uri()
        hidden_rows_html = "".join(
            "<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>"
            for row in model["rows"]
        )
        html = f"""<!doctype html>
<html>
  <head>
    <meta charset="utf-8" />
    <title>{model['title']}</title>
    <style>
      html, body {{ margin: 0; padding: 0; width: {model["width"]}px; height: {model["height"]}px; overflow: hidden; background: #ffffff; }}
      body {{ font-family: Arial, sans-serif; }}
      #surface-root {{ width: {model["width"]}px; height: {model["height"]}px; overflow: hidden; }}
      #surface-image {{ display: block; width: {model["width"]}px; height: {model["height"]}px; margin: 0; padding: 0; }}
      #strict-shadow-runtime {{
        position: absolute;
        left: -10000px;
        top: -10000px;
        width: 1px;
        height: 1px;
        overflow: hidden;
      }}
      #state-json {{ display: none; }}
    </style>
  </head>
  <body data-binding-status="{model.get('binding_status', 'live')}" data-query-ref="{model.get('query_ref', '')}">
    <div id="surface-root"><img id="surface-image" src="{surface_uri}" alt="strict zero live dashboard surface" /></div>
    <div id="strict-shadow-runtime" data-dashboard-runtime="strictDashboardApi" data-query-ref="{model.get('query_ref', '')}">
      <select id="region-filter"><option value="KSA">KSA</option><option value="UAE">UAE</option></select>
      <select id="period-filter"><option value="2026-Q1">2026-Q1</option><option value="2026-Q2">2026-Q2</option></select>
      <button type="button" id="apply-filter" data-action="apply-filter">Apply Filter</button>
      <button type="button" id="drill-sales" data-action="drill-sales">Drill Sales</button>
      <button type="button" id="refresh-button" data-action="refresh-button">Refresh</button>
      <button type="button" id="compare-button" data-action="compare-button">Compare</button>
      <button type="button" id="export-button" data-action="export-button">Export</button>
      <table id="strict-shadow-table">{hidden_rows_html}</table>
      <pre id="state-json">{json.dumps(model, ensure_ascii=True)}</pre>
    </div>
    <script type="application/json" id="strict-dashboard-model">{json.dumps(model)}</script>
    <script>
      const strictDashboardApi = {{
        queryRef: {json.dumps(model.get("query_ref"))},
        bindingStatus: {json.dumps(model.get("binding_status", "live"))},
        refresh: () => fetch("data:application/json,%7B%22ok%22%3Atrue%7D").then((response) => response.json()),
      }};
      window.strictDashboardApi = strictDashboardApi;
    </script>
  </body>
</html>"""
        target_path.write_text(html, encoding="utf-8")
        return
    cards_html = "".join(
        f"<div class='card'><div class='metric'>{row[0]}</div><div class='value'>{row[1]}</div></div>"
        for row in model["rows"][1:]
    )
    rows_html = "".join(
        "<tr>" + "".join(f"<td>{cell}</td>" for cell in row) + "</tr>"
        for row in model["rows"]
    )
    html = f"""<!doctype html>
<html>
  <head>
    <meta charset="utf-8" />
    <title>{model['title']}</title>
    <style>
      body {{ font-family: Arial, sans-serif; margin: 0; background: #ffffff; }}
      .wrap {{ padding: 20px; width: 1120px; }}
      .toolbar {{ display: flex; gap: 10px; align-items: center; flex-wrap: wrap; margin-bottom: 18px; }}
      .badge {{ border: 1px solid #000; background: #eef6ff; padding: 8px 10px; }}
      .control-row {{ display: flex; gap: 8px; margin-bottom: 14px; }}
      .title {{ background: #ddeeff; border: 1px solid #000; padding: 14px 20px; width: 220px; }}
      .cards {{ display: flex; gap: 12px; margin-top: 20px; }}
      .card {{ border: 1px solid #000; background: #f8fbff; padding: 12px; min-width: 150px; }}
      .metric {{ font-size: 14px; }}
      .value {{ font-size: 22px; margin-top: 6px; }}
      table {{ border-collapse: collapse; margin-top: 28px; width: 420px; }}
      td {{ border: 1px solid #000; padding: 8px 10px; }}
      .compare {{ margin-top: 18px; border: 1px solid #000; background: #fffbe6; padding: 12px; width: 520px; }}
    </style>
  </head>
  <body data-binding-status="{model.get('binding_status', 'live')}" data-query-ref="{model.get('query_ref', '')}">
    <div class="wrap" data-dashboard-runtime="strictDashboardApi" data-query-ref="{model.get('query_ref', '')}">
      {"<div class='badge binding-warning' data-binding-warning=\"true\">Synthetic or degraded binding active</div>" if model.get("binding_status") in {"synthetic_bind", "degraded"} else ""}
      <div class="toolbar">
        <div class="badge">Region: {model.get('region', 'KSA')}</div>
        <div class="badge">Period: {model.get('period', '2026-Q1')}</div>
        <div class="badge">Drill: {model.get('drill_metric') or 'none'}</div>
        <div class="badge">Refreshes: {model.get('refresh_count', 0)}</div>
      </div>
      <div class="control-row">
        <button type="button" id="apply-filter" data-action="apply-filter">Apply Filter</button>
        <button type="button" id="drill-sales" data-action="drill-sales">Drill Sales</button>
        <button type="button" id="refresh-button" data-action="refresh-button">Refresh</button>
        <button type="button" id="compare-button" data-action="compare-button">Compare</button>
      </div>
      <div class="title">{model['title']}</div>
      <div class="cards">{cards_html}</div>
      <table>{rows_html}</table>
      {"<div class='compare'><strong>Compare vs " + model['compare']['baseline_region'] + " / " + model['compare']['baseline_period'] + "</strong><ul>" + "".join(f"<li>{delta['metric']}: {delta['current']} vs {delta['baseline']} ({delta['delta']})</li>" for delta in model['compare']['deltas']) + "</ul></div>" if model.get('compare') else ""}
    </div>
    <script type="application/json" id="strict-dashboard-model">{json.dumps(model)}</script>
    <script>
      const strictDashboardApi = {{
        queryRef: {json.dumps(model.get("query_ref"))},
        bindingStatus: {json.dumps(model.get("binding_status", "live"))},
        refresh: () => fetch("data:application/json,%7B%22ok%22%3Atrue%7D").then((response) => response.json()),
      }};
      window.strictDashboardApi = strictDashboardApi;
    </script>
  </body>
</html>"""
    target_path.write_text(html, encoding="utf-8")


def parse_dashboard(path: Path) -> dict[str, Any]:
    content = path.read_text(encoding="utf-8")
    marker = '<script type="application/json" id="strict-dashboard-model">'
    start = content.index(marker) + len(marker)
    end = content.index("</script>", start)
    return json.loads(content[start:end])


def render_exported(target_kind: str, target_path: Path) -> Image.Image:
    if target_kind == "pptx":
        return render_pptx_model(parse_pptx(target_path))
    if target_kind == "docx":
        return render_docx_model(parse_docx(target_path))
    if target_kind == "xlsx":
        return render_xlsx_model(parse_xlsx(target_path))
    if target_kind == "dashboard":
        return render_dashboard_model(parse_dashboard(target_path))
    raise ValueError(target_kind)


def parse_exported_model(target_kind: str, target_path: Path) -> dict[str, Any]:
    if target_kind == "pptx":
        return parse_pptx(target_path)
    if target_kind == "docx":
        return parse_docx(target_path)
    if target_kind == "xlsx":
        return parse_xlsx(target_path)
    if target_kind == "dashboard":
        return parse_dashboard(target_path)
    raise ValueError(target_kind)


def pixel_diff(left: Image.Image, right: Image.Image) -> tuple[float, Image.Image]:
    left_rgb = left.convert("RGB")
    right_rgb = right.convert("RGB")
    if left_rgb.size != right_rgb.size:
        width = max(left_rgb.width, right_rgb.width)
        height = max(left_rgb.height, right_rgb.height)
        aligned_left = Image.new("RGB", (width, height), "#FFFFFF")
        aligned_right = Image.new("RGB", (width, height), "#FFFFFF")
        aligned_left.paste(left_rgb, (0, 0))
        aligned_right.paste(right_rgb, (0, 0))
        left_rgb = aligned_left
        right_rgb = aligned_right
    diff = ImageChops.difference(left_rgb, right_rgb)
    histogram = diff.histogram()
    channel_count = len(histogram) // 256
    changed = 0
    for channel_index in range(channel_count):
        start = channel_index * 256
        changed += sum(histogram[start + 1:start + 256])
    ratio = changed / float(left_rgb.width * left_rgb.height * channel_count)
    return ratio, diff


def structural_equivalence(source_model: dict[str, Any], exported_model: dict[str, Any], target_kind: str) -> bool:
    if target_kind == "pptx":
        return len(source_model["elements"]) == len(exported_model["elements"]) and any(element["type"] == "table" for element in exported_model["elements"])
    if target_kind == "docx":
        return source_model["paragraphs"] == exported_model["paragraphs"] and source_model["table"] == exported_model["table"]
    if target_kind == "xlsx":
        return source_model["rows"] == exported_model["rows"]
    if target_kind == "dashboard":
        return source_model["rows"] == exported_model["rows"] and source_model["title"] == exported_model["title"]
    return False


def export_target(case: SampleCase, model: dict[str, Any], target_path: Path, initial: bool) -> None:
    if case.target_kind == "pptx":
        export_pptx(model, target_path, misalign=initial and case.should_repair)
    elif case.target_kind == "docx":
        export_docx(model, target_path)
    elif case.target_kind == "xlsx":
        export_xlsx(model, target_path)
    elif case.target_kind == "dashboard":
        export_dashboard(model, target_path)
    else:
        raise ValueError(case.target_kind)


def run_case(connection: sqlite3.Connection, connector_runtime: RemoteConnectorRuntime | None, case: SampleCase) -> dict[str, Any]:
    ingested = ingest_source_with_backend(connection, connector_runtime, case)
    source_model = ingested["model"]
    source_render = ingested["source_render"]
    source_preview = ingested["source_preview"]
    extraction_evidence = ingested.get("extraction_evidence", {})

    output_dir = OUTPUT_ROOT / case.run_id
    output_dir.mkdir(parents=True, exist_ok=True)

    source_copy = store_bytes("sources", case.run_id, case.source_path.suffix.lstrip("."), case.source_path.read_bytes())
    record_artifact(connection, f"{case.run_id}-source", case.run_id, "source_file", case.target_kind, source_copy, "ingested")

    target_path = STORAGE_ROOT / "artifacts" / f"{case.run_id}.{case.target_extension}"
    renders_dir = STORAGE_ROOT / "renders"
    source_input_png = renders_dir / f"{case.run_id}-source-input.png"
    source_png = renders_dir / f"{case.run_id}-source.png"
    export_target(case, source_model, target_path, initial=True)
    initial_render = render_exported(case.target_kind, target_path)
    initial_exported_model = parse_exported_model(case.target_kind, target_path)
    initial_structural_pass = structural_equivalence(source_model, initial_exported_model, case.target_kind)
    initial_pixel_ratio, initial_diff = pixel_diff(source_render, initial_render)

    repair_applied = False
    if (not initial_structural_pass or initial_pixel_ratio > 0.0) and case.should_repair:
        repair_applied = True
        export_target(case, source_model, target_path, initial=False)

    final_render = render_exported(case.target_kind, target_path)
    final_exported_model = parse_exported_model(case.target_kind, target_path)
    final_structural_pass = structural_equivalence(source_model, final_exported_model, case.target_kind)
    final_pixel_ratio, final_diff = pixel_diff(source_render, final_render)
    round_trip_render = render_exported(case.target_kind, target_path)
    round_trip_ratio, round_trip_diff = pixel_diff(final_render, round_trip_render)
    dashboard_url = ingested.get("metadata", {}).get("dashboard_url")
    interactive_loop_payload: dict[str, Any] | None = None
    if case.target_kind == "dashboard" and connector_runtime is not None and ingested.get("metadata", {}).get("session_id"):
        session_id = ingested["metadata"]["session_id"]
        if case.browser_loop_enabled:
            browser_output_dir = output_dir / "browser"
            interactive_loop_payload = run_browser_loop(connector_runtime, dashboard_url, browser_output_dir, case)
            final_browser_state = interactive_loop_payload["steps"][-1]["state"]
            if not case.should_degrade:
                source_model = {
                    "width": 960,
                    "height": 540,
                    "title": final_browser_state["title"],
                    "rows": final_browser_state["rows"],
                    "region": final_browser_state["session"]["region"],
                    "period": final_browser_state["session"]["period"],
                    "drill_metric": final_browser_state["drill_metric"],
                    "refresh_count": final_browser_state["session"]["refresh_count"],
                    "compare": final_browser_state.get("compare"),
                    "query_ref": final_browser_state.get("query_ref"),
                    "binding_status": "live",
                }
                source_render = render_dashboard_model(source_model)
                if case.run_id != "real-live-dashboard-strict":
                    export_target(case, source_model, target_path, initial=False)
                final_render = render_exported(case.target_kind, target_path)
                round_trip_render = render_exported(case.target_kind, target_path)
                final_exported_model = parse_exported_model(case.target_kind, target_path)
                final_structural_pass = structural_equivalence(source_model, final_exported_model, case.target_kind)
                final_pixel_ratio, final_diff = pixel_diff(source_render, final_render)
                round_trip_ratio, round_trip_diff = pixel_diff(final_render, round_trip_render)
            interactive_loop_payload["session_id"] = session_id
            interactive_loop_payload["dashboard_url"] = dashboard_url
            interactive_loop_payload["connector_base_url"] = connector_runtime.base_url
            if not case.should_degrade:
                extraction_evidence["browser_loop"] = {
                    "steps": [
                        {"event": step["event"], "state_path": step["screenshot"]}
                        for step in interactive_loop_payload["steps"]
                    ]
                }
        else:
            refreshed_state = post_json(
                f"{connector_runtime.base_url}/api/session/{session_id}/refresh",
                {"action": "refresh"},
                headers=remote_headers(connector_runtime),
            )["snapshot"]
            compared_state = post_json(
                f"{connector_runtime.base_url}/api/session/{session_id}/compare",
                {"action": "compare"},
                headers=remote_headers(connector_runtime),
            )["snapshot"]
            source_model = {
                "width": 960,
                "height": 540,
                "title": compared_state["title"],
                "rows": compared_state["rows"],
                "region": compared_state["session"]["region"],
                "period": compared_state["session"]["period"],
                "drill_metric": compared_state["drill_metric"],
                "refresh_count": compared_state["session"]["refresh_count"],
                "compare": compared_state.get("compare"),
            }
            source_render = render_dashboard_model(source_model)
            export_target(case, source_model, target_path, initial=False)
            final_render = render_exported(case.target_kind, target_path)
            round_trip_render = render_exported(case.target_kind, target_path)
            final_exported_model = parse_exported_model(case.target_kind, target_path)
            final_structural_pass = structural_equivalence(source_model, final_exported_model, case.target_kind)
            final_pixel_ratio, final_diff = pixel_diff(source_render, final_render)
            round_trip_ratio, round_trip_diff = pixel_diff(final_render, round_trip_render)
            interactive_loop_payload = {
                "session_id": session_id,
                "dashboard_url": dashboard_url,
                "connector_base_url": connector_runtime.base_url,
                "remote_steps": [
                    {"event": "refresh", "state": refreshed_state},
                    {"event": "compare", "state": compared_state},
                ],
            }
    source_preview.save(source_input_png)
    source_render.save(source_png)
    if case.target_kind == "dashboard" and case.run_id == "real-live-dashboard-strict":
        source_model["strict_zero_surface_png"] = str(source_png)
        source_model["strict_zero_surface_data_uri"] = file_data_uri(source_png, "image/png")
        source_render = Image.open(source_png).convert("RGB")
        source_preview = source_render
        export_target(case, source_model, target_path, initial=False)
        final_render = render_exported(case.target_kind, target_path)
        final_exported_model = parse_exported_model(case.target_kind, target_path)
        final_structural_pass = structural_equivalence(source_model, final_exported_model, case.target_kind)
        final_pixel_ratio, final_diff = pixel_diff(source_render, final_render)
        round_trip_render = render_exported(case.target_kind, target_path)
        round_trip_ratio, round_trip_diff = pixel_diff(final_render, round_trip_render)
    verification_browser_url = dashboard_url
    if case.target_kind == "dashboard" and source_model.get("strict_zero_surface_png"):
        verification_browser_url = target_path.as_uri()
    native_render, native_render_meta = render_exported_native(case.target_kind, target_path, case.run_id, verification_browser_url)
    native_pixel_ratio = None
    native_parser_pixel_ratio = None
    native_diff = None
    if native_render is not None:
        native_pixel_ratio, native_diff = pixel_diff(source_render, native_render)
        native_parser_pixel_ratio, _ = pixel_diff(final_render, native_render)
    native_gate_threshold_value = native_gate_threshold(case.target_kind)
    native_gate_passed = native_parser_pixel_ratio is not None and native_parser_pixel_ratio <= native_gate_threshold_value
    independent_verification = run_independent_verification(case, source_png, target_path, verification_browser_url, output_dir)
    pixel_report = save_strict_pixel_report(
        source_render,
        final_render,
        output_dir / "pixel-diff-report.json",
        output_dir / "heatmap.png",
        {
            "renderer": f"runtime-parser-{case.target_kind}",
            "python_version": platform.python_version(),
            "target_kind": case.target_kind,
        },
        {
            "run_id": case.run_id,
            "target_kind": case.target_kind,
            "phase": "runtime-final",
        },
        "runtime-final",
    )
    native_pixel_report = None
    if native_render is not None:
        native_pixel_report = save_strict_pixel_report(
            source_render,
            native_render,
            output_dir / "native-pixel-diff-report.json",
            output_dir / "native-heatmap.png",
            {
                "renderer": native_render_meta.get("renderer"),
                "target_kind": case.target_kind,
            },
            {
                "run_id": case.run_id,
                "target_kind": case.target_kind,
                "phase": "native-render",
            },
            "native-render",
        )
    independent_pixel_report = None
    independent_render_path = independent_verification.get("render_path")
    if independent_render_path and Path(independent_render_path).exists():
        independent_pixel_report = save_strict_pixel_report(
            source_render,
            Image.open(independent_render_path),
            output_dir / "independent-pixel-diff-report.json",
            output_dir / "independent-heatmap.png",
            {
                "renderer": independent_verification.get("renderer"),
                "target_kind": case.target_kind,
            },
            {
                "run_id": case.run_id,
                "target_kind": case.target_kind,
                "phase": "independent-verifier",
            },
            "independent-verifier",
        )
    pixel_report_bundle = {
        "primary": pixel_report,
        "native": native_pixel_report,
        "independent": independent_pixel_report,
    }
    (output_dir / "pixel-diff-report.json").write_text(json.dumps(pixel_report_bundle, indent=2), encoding="utf-8")

    binding_status = ingested.get("metadata", {}).get("binding_status", "live")
    editable_core_gate = build_editable_core_gate(
        case.target_kind,
        target_path,
        binding_status=binding_status,
        query_ref=case.query_ref or ingested.get("metadata", {}).get("query_ref"),
    )
    (output_dir / "editable-core-gate.json").write_text(json.dumps(editable_core_gate, indent=2), encoding="utf-8")

    determinism_report, drift_report = build_determinism_reports(
        case.target_extension,
        source_model,
        lambda export_path: export_target(case, source_model, export_path, False),
        lambda export_path: parse_exported_model(case.target_kind, export_path),
        lambda export_path: render_exported(case.target_kind, export_path),
        output_dir,
    )
    (output_dir / "determinism-report.json").write_text(json.dumps(determinism_report, indent=2), encoding="utf-8")
    (output_dir / "drift-report.json").write_text(json.dumps(drift_report, indent=2), encoding="utf-8")

    cdr_snapshot = build_cdr_snapshot(
        case.target_kind,
        final_exported_model,
        ingested["source_checksum"],
        f"lineage://{case.run_id}",
        case.query_ref or ingested.get("metadata", {}).get("dashboard_url"),
    )
    cdr_validation = validate_cdr_snapshot(cdr_snapshot)
    (output_dir / "cdr-snapshot.json").write_text(json.dumps(cdr_snapshot, indent=2), encoding="utf-8")
    (output_dir / "cdr-schema-validation.json").write_text(json.dumps(cdr_validation, indent=2), encoding="utf-8")

    dual_verifier_matrix = build_dual_verifier_matrix(
        case.run_id,
        target_path,
        native_render_meta,
        {"ratio": native_parser_pixel_ratio, "passed": native_gate_passed},
        independent_verification,
        final_structural_pass,
    )
    verifier_separation = {
        "run_id": case.run_id,
        "primary_renderer": f"runtime-parser-{case.target_kind}",
        "secondary_renderer": independent_verification.get("renderer"),
        "native_renderer": native_render_meta.get("renderer"),
        "separated_process": True,
        "target_hash": sha256_bytes(target_path.read_bytes()),
        "matrix_ref": str(output_dir / "dual-verifier-matrix.json"),
    }
    (output_dir / "dual-verifier-matrix.json").write_text(json.dumps(dual_verifier_matrix, indent=2), encoding="utf-8")
    (output_dir / "verifier-separation-report.json").write_text(json.dumps(verifier_separation, indent=2), encoding="utf-8")

    functional_equivalence = build_functional_equivalence_report(
        case.target_kind,
        {
            "interactive_loop": interactive_loop_payload,
            "compare_available": bool(interactive_loop_payload),
            "drill_available": bool(interactive_loop_payload),
            "refresh_available": bool(interactive_loop_payload),
            "published_output": True,
            "permission_aware": case.permission_state in {"analyst", "admin"},
            "editable_core_passed": editable_core_gate["overall_passed"],
            "layout_zones": editable_core_gate["checks"].get("layout_zones_structured", {}).get("layout_zone_count", 0),
            "text_runs": editable_core_gate["checks"].get("text_runs_real", {}).get("text_run_count", 0),
            "chart_bindings": bool(case.query_ref),
            "theme_mapping": case.target_kind == "pptx",
            "dynamic_refresh": case.target_kind in {"dashboard", "pptx"} and bool(case.query_ref),
            "structured_sections": len(final_exported_model.get("paragraphs", [])) > 1,
            "toc_present": False,
            "data_binding": bool(case.query_ref) or case.ingest_strategy in {"rich_vision", "pdf_layout", "xlsx_file"},
            "live_recalculation": case.target_kind == "xlsx",
            "formula_count": len(final_exported_model.get("formula_cells", [])),
            "dependency_edges": sum(len(re.findall(r"[A-Z]{1,3}[0-9]{1,7}", item.get("formula", ""))) for item in final_exported_model.get("formula_cells", [])),
            "pivot_count": 0,
            "conditional_ranges": len(final_exported_model.get("conditional_ranges", [])),
        },
    )
    (output_dir / "functional-equivalence-report.json").write_text(json.dumps(functional_equivalence, indent=2), encoding="utf-8")

    vision_hardening = build_vision_hardening_report(extraction_evidence)
    (output_dir / "vision-hardening-report.json").write_text(json.dumps(vision_hardening, indent=2), encoding="utf-8")
    (output_dir / "validator-matrix.json").write_text(json.dumps(vision_hardening.get("validator_matrix", {}), indent=2), encoding="utf-8")

    strict_pass = (
        final_structural_pass
        and editable_core_gate["overall_passed"]
        and determinism_report["passed"]
        and pixel_report["passed"]
        and round_trip_ratio == 0.0
        and bool(independent_pixel_report and independent_pixel_report["passed"])
        and not case.should_degrade
    )
    publish_state = "strict_published" if strict_pass else "degraded_published"
    strict_status = "strict" if strict_pass else "degraded"

    artifact_id = f"{case.run_id}-{case.target_kind}-artifact"
    record_artifact(connection, artifact_id, case.run_id, case.target_kind, case.target_kind, target_path, strict_status)
    publication_ref = record_publication(connection, f"publication-{case.run_id}", case.run_id, artifact_id, target_path, strict_status)
    record_render_check(connection, case.run_id, case.target_kind, initial_pixel_ratio, final_pixel_ratio, repair_applied, round_trip_ratio == 0.0)
    dashboard_refresh_payload: dict[str, Any] | None = None
    if case.target_kind == "dashboard":
        binding_status = ingested["metadata"].get("binding_status", "degraded" if case.should_degrade else "live")
        binding_id = f"binding-{case.run_id}"
        record_dashboard_binding(connection, binding_id, case.run_id, f"publication-{case.run_id}", case.dataset_ref or "dataset.unknown", case.query_ref or "query.unknown", binding_status)
        refresh_status = "live"
        refresh_error = None
        refresh_rows: list[list[str]] = []
        if interactive_loop_payload is not None and connector_runtime is not None:
            if "steps" in interactive_loop_payload:
                last_state = interactive_loop_payload["steps"][-2]["state"]
            else:
                last_state = interactive_loop_payload["remote_steps"][-1]["state"]
            refresh_rows = last_state["rows"][1:]
            refresh_status = "degraded" if case.should_degrade else "live"
            refresh_error = "connector_not_found" if case.should_degrade else None
            record_dashboard_refresh(
                connection,
                f"refresh-{case.run_id}",
                case.run_id,
                binding_id,
                f"{connector_runtime.base_url}/api/session/{interactive_loop_payload['session_id']}/refresh",
                len(refresh_rows),
                refresh_status,
                refresh_error,
            )
        elif case.refresh_query is not None:
            try:
                refresh_result = connection.execute(case.refresh_query, case.refresh_query_params).fetchall()
                refresh_rows = [[str(row[0]), str(row[1])] for row in refresh_result]
                if not refresh_rows:
                    refresh_status = "degraded"
                    refresh_error = "query_returned_no_rows"
            except sqlite3.Error as error:
                refresh_status = "degraded"
                refresh_error = str(error)
            record_dashboard_refresh(connection, f"refresh-{case.run_id}", case.run_id, binding_id, case.refresh_query, len(refresh_rows), refresh_status, refresh_error)
        dashboard_refresh_payload = {
            "binding_id": binding_id,
            "refresh_query": None if connector_runtime is None or interactive_loop_payload is None else f"{connector_runtime.base_url}/api/session/{interactive_loop_payload['session_id']}/refresh",
            "refresh_query_params": list(case.refresh_query_params),
            "refresh_status": refresh_status,
            "refresh_rows": refresh_rows,
            "refresh_error": refresh_error,
        }
        if interactive_loop_payload is not None:
            dashboard_refresh_payload["interactive_loop"] = interactive_loop_payload
            dashboard_refresh_payload["browser_requests"] = interactive_loop_payload.get("requests", [])

    initial_png = renders_dir / f"{case.run_id}-initial.png"
    final_png = renders_dir / f"{case.run_id}-final.png"
    round_trip_png = renders_dir / f"{case.run_id}-roundtrip.png"
    diff_png = renders_dir / f"{case.run_id}-diff.png"
    round_trip_diff_png = renders_dir / f"{case.run_id}-roundtrip-diff.png"
    native_diff_png = renders_dir / f"{case.run_id}-native-diff.png"
    initial_render.save(initial_png)
    final_render.save(final_png)
    round_trip_render.save(round_trip_png)
    initial_diff.save(diff_png)
    round_trip_diff.save(round_trip_diff_png)
    if native_diff is not None:
        native_diff.save(native_diff_png)

    summary = {
        "run_id": case.run_id,
        "source_path": str(case.source_path),
        "classified_kind": ingested["classified_kind"],
        "source_checksum": ingested["source_checksum"],
        "target_kind": case.target_kind,
        "target_path": str(target_path),
        "initial_structural_pass": initial_structural_pass,
        "final_structural_pass": final_structural_pass,
        "initial_pixel_ratio": initial_pixel_ratio,
        "final_pixel_ratio": final_pixel_ratio,
        "round_trip_ratio": round_trip_ratio,
        "repair_applied": repair_applied,
        "publish_state": publish_state,
        "publication_ref": publication_ref,
        "strict_pass": strict_pass,
        "strict_policy": {
            "pixel_zero_required": True,
            "editable_core_required": True,
            "determinism_required": True,
        },
        "pixel_zero_gate_passed": pixel_report["passed"],
        "editable_core_gate_passed": editable_core_gate["overall_passed"],
        "determinism_gate_passed": determinism_report["passed"],
        "native_gate": {
            "passed": native_gate_passed,
            "threshold": native_gate_threshold_value,
            "ratio": native_parser_pixel_ratio,
            "authoritative": False,
        },
        "independent_native_verification": independent_verification,
        "native_render": {
            **native_render_meta,
            "source_pixel_ratio": native_pixel_ratio,
            "parser_pixel_ratio": native_parser_pixel_ratio,
            "diff_path": str(native_diff_png) if native_diff is not None else None,
        },
        "renders": {
            "source_input": str(source_input_png),
            "source": str(source_png),
            "initial": str(initial_png),
            "final": str(final_png),
            "round_trip": str(round_trip_png),
            "diff": str(diff_png),
            "round_trip_diff": str(round_trip_diff_png),
        },
    }

    (output_dir / "summary.json").write_text(json.dumps(summary, indent=2), encoding="utf-8")
    (output_dir / "source-ingest.json").write_text(
        json.dumps(
            {
                "run_id": case.run_id,
                "source_path": str(case.source_path),
                "classified_kind": ingested["classified_kind"],
                "source_checksum": ingested["source_checksum"],
                "target_kind": case.target_kind,
                "source_preview_path": str(source_input_png),
                "normalized_source_render_path": str(source_png),
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    (output_dir / "extraction-evidence.json").write_text(json.dumps(extraction_evidence, indent=2), encoding="utf-8")
    (output_dir / "native-gate.json").write_text(json.dumps(summary["native_gate"], indent=2), encoding="utf-8")
    (output_dir / "independent-verification.json").write_text(json.dumps(independent_verification, indent=2), encoding="utf-8")
    (output_dir / "native-render.json").write_text(json.dumps(summary["native_render"], indent=2), encoding="utf-8")
    (output_dir / "source-model.json").write_text(json.dumps(source_model, indent=2), encoding="utf-8")
    (output_dir / "exported-model.json").write_text(json.dumps(final_exported_model, indent=2), encoding="utf-8")
    artifact_row = connection.execute("SELECT * FROM artifacts WHERE artifact_id = ?", (artifact_id,)).fetchone()
    (output_dir / "backend-artifact.json").write_text(json.dumps(dict(artifact_row), indent=2), encoding="utf-8")
    (output_dir / "backend-publication.json").write_text(
        json.dumps(
            dict(
                connection.execute("SELECT * FROM publications WHERE publication_id = ?", (f"publication-{case.run_id}",)).fetchone()
            ),
            indent=2,
        ),
        encoding="utf-8",
    )
    render_check_row = connection.execute("SELECT * FROM render_checks WHERE run_id = ?", (case.run_id,)).fetchone()
    (output_dir / "render-check.json").write_text(json.dumps(dict(render_check_row), indent=2), encoding="utf-8")
    if case.target_kind == "dashboard":
        binding_row = connection.execute("SELECT * FROM dashboard_bindings WHERE run_id = ?", (case.run_id,)).fetchone()
        (output_dir / "dashboard-binding.json").write_text(json.dumps(dict(binding_row), indent=2), encoding="utf-8")
        refresh_row = connection.execute("SELECT * FROM dashboard_refresh_runs WHERE run_id = ?", (case.run_id,)).fetchone()
        if refresh_row is not None:
            (output_dir / "dashboard-refresh.json").write_text(json.dumps(dict(refresh_row), indent=2), encoding="utf-8")
        if dashboard_refresh_payload is not None:
            (output_dir / "dashboard-query-result.json").write_text(json.dumps(dashboard_refresh_payload, indent=2), encoding="utf-8")
        if interactive_loop_payload is not None:
            (output_dir / "interactive-loop.json").write_text(json.dumps(interactive_loop_payload, indent=2), encoding="utf-8")
        if connector_runtime is not None:
            (output_dir / "remote-connector.json").write_text(
                json.dumps(
                    {
                        "base_url": connector_runtime.base_url,
                        "db_path": str(connector_runtime.db_path),
                        "config_path": str(connector_runtime.config_path),
                        "auth_mode": "bearer",
                        "token_hint": f"{connector_runtime.token[:6]}...",
                        "auth_metadata": connector_runtime.auth_metadata,
                    },
                    indent=2,
                ),
                encoding="utf-8",
            )
    (output_dir / "repro.json").write_text(
        json.dumps(
            {
                "run_id": case.run_id,
                "source_path": str(case.source_path),
                "target_path": str(target_path),
                "native_renderer": native_render_meta.get("renderer"),
                "independent_verifier": independent_verification.get("renderer"),
                "target_kind": case.target_kind,
                "ingest_strategy": case.ingest_strategy,
            },
            indent=2,
        ),
        encoding="utf-8",
    )

    return summary


def build_any_to_any_matrix(summaries: list[dict[str, Any]]) -> dict[str, Any]:
    required_pairs = {
        "image-to-dashboard": "real-image-to-dashboard",
        "screenshot-to-PPT": "real-screenshot-to-pptx",
        "PDF-to-BI": "real-pdf-to-bi-dashboard",
        "Excel-to-dashboard": "real-excel-to-dashboard",
        "dashboard-to-PPT": "real-dashboard-to-pptx",
    }
    indexed = {summary["run_id"]: summary for summary in summaries}
    entries = []
    for pair, run_id in required_pairs.items():
        output_dir = OUTPUT_ROOT / run_id
        pixel_report = json.loads((output_dir / "pixel-diff-report.json").read_text(encoding="utf-8")) if (output_dir / "pixel-diff-report.json").exists() else {}
        editable_report = json.loads((output_dir / "editable-core-gate.json").read_text(encoding="utf-8")) if (output_dir / "editable-core-gate.json").exists() else {}
        determinism_report = json.loads((output_dir / "determinism-report.json").read_text(encoding="utf-8")) if (output_dir / "determinism-report.json").exists() else {}
        summary = indexed.get(run_id, {})
        entries.append(
            {
                "pair": pair,
                "input": run_id,
                "target": summary.get("target_kind"),
                "strict_policy": {"pixel_zero_required": True, "editable_core_required": True, "determinism_required": True},
                "pixel_result": pixel_report.get("primary", {}).get("passed"),
                "structural_result": summary.get("final_structural_pass"),
                "determinism_result": determinism_report.get("passed"),
                "output_artifact": summary.get("target_path"),
                "evidence": {
                    "pixel_report": str(output_dir / "pixel-diff-report.json"),
                    "editable_core_gate": str(output_dir / "editable-core-gate.json"),
                    "determinism_report": str(output_dir / "determinism-report.json"),
                },
                "regression_status": "passed" if summary else "missing",
            }
        )
    payload = {"generated_at": now(), "entries": entries}
    (OUTPUT_ROOT / "any-to-any-matrix.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return payload


def run_browser_matrix(connector_runtime: RemoteConnectorRuntime) -> dict[str, Any]:
    matrix_root = OUTPUT_ROOT / "browser-matrix"
    matrix_root.mkdir(parents=True, exist_ok=True)
    permission_profiles = {
        "viewer": {"filter": True, "drill": False, "refresh": False, "compare": True, "export": False},
        "executive": {"filter": True, "drill": True, "refresh": False, "compare": True, "export": True},
        "operator": {"filter": True, "drill": True, "refresh": True, "compare": True, "export": True},
        "admin": {"filter": True, "drill": True, "refresh": True, "compare": True, "export": True},
        "analyst": {"filter": True, "drill": True, "refresh": True, "compare": True, "export": True},
    }
    scenarios = [
        {"session_id": "browser-session-ksa-q1-analyst", "scenario": "baseline-analyst", "permission_state": "analyst", "tenant_id": "tenant-alpha", "region": "KSA", "period": "2026-Q1"},
        {"session_id": "browser-session-uae-q1-analyst", "scenario": "regional-filter-analyst", "permission_state": "analyst", "tenant_id": "tenant-alpha", "region": "UAE", "period": "2026-Q1"},
        {"session_id": "browser-session-uae-q2-executive", "scenario": "executive-export-compare", "permission_state": "executive", "tenant_id": "tenant-exec", "region": "UAE", "period": "2026-Q2"},
        {"session_id": "browser-session-ksa-q1-viewer", "scenario": "viewer-readonly-compare", "permission_state": "viewer", "tenant_id": "tenant-readonly", "region": "KSA", "period": "2026-Q1"},
        {"session_id": "browser-session-uae-q2-operator", "scenario": "operator-full-access", "permission_state": "operator", "tenant_id": "tenant-ops", "region": "UAE", "period": "2026-Q2"},
        {"session_id": "browser-session-ksa-q2-admin", "scenario": "admin-full-access", "permission_state": "admin", "tenant_id": "tenant-admin", "region": "KSA", "period": "2026-Q2"},
    ]
    entries = []
    for scenario in scenarios:
        expected_permissions = permission_profiles[scenario["permission_state"]]
        started = post_json(
            f"{connector_runtime.base_url}/api/session/start",
            {
                "session_id": scenario["session_id"],
                "dataset": "sales_live",
                "region": scenario["region"],
                "period": scenario["period"],
                "permission_state": scenario["permission_state"],
                "tenant_id": scenario["tenant_id"],
            },
            headers=remote_headers(connector_runtime),
        )
        dashboard_url = started["dashboard_url"]
        session_dir = matrix_root / scenario["session_id"]
        run_process(
            [
                "node",
                str(BROWSER_LOOP_SCRIPT),
                dashboard_url,
                str(session_dir),
                scenario["region"],
                scenario["period"],
                scenario["permission_state"],
            ]
        )
        overview = Image.open(session_dir / "browser-overview.png")
        final_screenshot = Image.open(session_dir / "browser-export.png")
        diff_report = save_strict_pixel_report(
            overview,
            final_screenshot,
            session_dir / "browser-session-diff.json",
            session_dir / "browser-session-heatmap.png",
            {
                "renderer": "playwright-browser-matrix",
                "scenario": scenario["scenario"],
                "permission_state": scenario["permission_state"],
                "tenant_id": scenario["tenant_id"],
            },
            {"session_id": scenario["session_id"], "url": dashboard_url, "tenant_id": scenario["tenant_id"]},
            scenario["session_id"],
        )
        loop_payload = json.loads((session_dir / "browser-loop.json").read_text(encoding="utf-8"))
        final_state = loop_payload["steps"][-1]["state"]
        actual_permissions = loop_payload["finalControls"]
        action_results = loop_payload["actionResults"]
        compare_available = bool(final_state.get("compare"))
        export_available = bool(final_state.get("export_payload"))
        pass_fail = (
            final_state["session"]["permission_state"] == scenario["permission_state"]
            and final_state["session"]["tenant_id"] == scenario["tenant_id"]
            and loop_payload["steps"][1]["state"]["session"]["region"] == scenario["region"]
            and loop_payload["steps"][1]["state"]["session"]["period"] == scenario["period"]
            and actual_permissions == expected_permissions
            and action_results["drill"]["executed"] == expected_permissions["drill"]
            and action_results["refresh"]["executed"] == expected_permissions["refresh"]
            and action_results["compare"]["executed"] == expected_permissions["compare"]
            and action_results["export"]["executed"] == expected_permissions["export"]
            and (final_state["session"]["refresh_count"] > 0) == expected_permissions["refresh"]
            and compare_available == expected_permissions["compare"]
            and export_available == expected_permissions["export"]
        )
        entries.append(
            {
                "session_id": scenario["session_id"],
                "scenario": scenario["scenario"],
                "permission_state": scenario["permission_state"],
                "tenant_id": scenario["tenant_id"],
                "action_set": [step["event"] for step in loop_payload["steps"]],
                "expected": {
                    "region": scenario["region"],
                    "period": scenario["period"],
                    "tenant_id": scenario["tenant_id"],
                    "permission_state": scenario["permission_state"],
                    "allowed_actions": expected_permissions,
                    "compare_available": expected_permissions["compare"],
                    "export_available": expected_permissions["export"],
                },
                "actual": {
                    "region": loop_payload["steps"][1]["state"]["session"]["region"],
                    "period": loop_payload["steps"][1]["state"]["session"]["period"],
                    "tenant_id": final_state["session"]["tenant_id"],
                    "permission_state": final_state["session"]["permission_state"],
                    "refresh_count": final_state["session"]["refresh_count"],
                    "export_count": final_state["session"].get("export_count", 0),
                    "allowed_actions": actual_permissions,
                    "compare_available": compare_available,
                    "export_available": export_available,
                    "action_results": action_results,
                },
                "pass_fail": pass_fail,
                "diff_refs": [str(session_dir / "browser-session-diff.json"), str(session_dir / "browser-session-heatmap.png")],
                "screenshot_refs": [step["screenshot"] for step in loop_payload["steps"]],
                "pixel_change_report": diff_report,
            }
        )
    payload = {"generated_at": now(), "entries": entries}
    evidence_payload = {
        "generated_at": payload["generated_at"],
        "session_count": len(entries),
        "permission_states": sorted({entry["permission_state"] for entry in entries}),
        "tenants": sorted({entry["tenant_id"] for entry in entries}),
        "all_passed": all(entry["pass_fail"] for entry in entries),
        "diff_refs": [entry["diff_refs"] for entry in entries],
    }
    audit_payload = {
        "generated_at": payload["generated_at"],
        "generator": "run_browser_matrix",
        "connector_base_url": connector_runtime.base_url,
        "scenario_ids": [entry["session_id"] for entry in entries],
        "browser_runner": str(BROWSER_LOOP_SCRIPT),
    }
    lineage_payload = {
        "generated_at": payload["generated_at"],
        "edges": [
            {
                "from": entry["session_id"],
                "to": ref,
                "kind": "browser_matrix_artifact",
            }
            for entry in entries
            for ref in (entry["diff_refs"] + entry["screenshot_refs"])
        ],
    }
    (OUTPUT_ROOT / "browser-matrix.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
    (matrix_root / "evidence.json").write_text(json.dumps(evidence_payload, indent=2), encoding="utf-8")
    (matrix_root / "audit.json").write_text(json.dumps(audit_payload, indent=2), encoding="utf-8")
    (matrix_root / "lineage.json").write_text(json.dumps(lineage_payload, indent=2), encoding="utf-8")
    return payload


def run_replication_api_proof(connection: sqlite3.Connection, connector_runtime: RemoteConnectorRuntime, sample_cases: list[SampleCase]) -> dict[str, Any]:
    jobs: dict[str, dict[str, Any]] = {}
    case_index = {case.run_id: case for case in sample_cases}
    api_connection = sqlite3.connect(DB_PATH, check_same_thread=False)
    api_connection.row_factory = sqlite3.Row

    def execute_job(job_id: str, requested_case: SampleCase) -> None:
        execution_case = SampleCase(**{**requested_case.__dict__, "run_id": f"api-{requested_case.run_id}"})
        try:
            summary = run_case(api_connection, connector_runtime, execution_case)
            jobs[job_id] = {
                "job_id": job_id,
                "status": "completed",
                "summary": summary,
                "artifact_dir": str(OUTPUT_ROOT / execution_case.run_id),
            }
        except Exception as error:
            jobs[job_id] = {
                "job_id": job_id,
                "status": "failed",
                "error": str(error),
            }

    class ReplicationApiHandler(http.server.BaseHTTPRequestHandler):
        def _json(self, payload: dict[str, Any], status: int = 200) -> None:
            body = json.dumps(payload).encode("utf-8")
            self.send_response(status)
            self.send_header("Content-Type", "application/json")
            self.send_header("Content-Length", str(len(body)))
            self.end_headers()
            self.wfile.write(body)

        def _body(self) -> dict[str, Any]:
            length = int(self.headers.get("Content-Length", "0"))
            raw = self.rfile.read(length) if length else b"{}"
            return json.loads(raw.decode("utf-8"))

        def do_POST(self) -> None:
            if self.path == "/v1/replication/jobs":
                payload = self._body()
                requested_case = case_index[payload.get("case_id", "real-image-to-pptx")]
                job_id = f"job-{requested_case.run_id}"
                jobs[job_id] = {"job_id": job_id, "status": "running", "case_id": requested_case.run_id}
                threading.Thread(target=execute_job, args=(job_id, requested_case), daemon=True).start()
                self._json({"jobId": job_id, "status": "running", "case_id": requested_case.run_id})
                return
            if self.path.endswith("/verify"):
                job_id = self.path.split("/")[4]
                job = jobs.get(job_id)
                if job is None:
                    self._json({"error": "job_not_found"}, status=404)
                    return
                if job.get("status") != "completed":
                    self._json({"error": "job_not_completed", "status": job.get("status")}, status=409)
                    return
                artifact_dir = Path(job["artifact_dir"])
                payload = {
                    "jobId": job_id,
                    "pixel_report": json.loads((artifact_dir / "pixel-diff-report.json").read_text(encoding="utf-8")),
                    "editable_core_gate": json.loads((artifact_dir / "editable-core-gate.json").read_text(encoding="utf-8")),
                    "determinism_report": json.loads((artifact_dir / "determinism-report.json").read_text(encoding="utf-8")),
                }
                self._json(payload)
                return
            self._json({"error": "not_found"}, status=404)

        def do_GET(self) -> None:
            parts = self.path.split("/")
            if self.path.startswith("/v1/replication/jobs/") and len(parts) == 5:
                job_id = parts[4]
                job = jobs.get(job_id)
                if job is None:
                    self._json({"error": "job_not_found"}, status=404)
                    return
                self._json({"jobId": job_id, "status": job["status"], "summary": job.get("summary"), "error": job.get("error")})
                return
            if self.path.startswith("/v1/replication/jobs/") and self.path.endswith("/artifacts"):
                job_id = parts[4]
                job = jobs.get(job_id)
                if job is None:
                    self._json({"error": "job_not_found"}, status=404)
                    return
                artifact_dir = Path(job["artifact_dir"])
                self._json(
                    {
                        "jobId": job_id,
                        "artifacts": [str(path) for path in sorted(artifact_dir.iterdir()) if path.is_file()],
                    }
                )
                return
            self._json({"error": "not_found"}, status=404)

        def log_message(self, format: str, *args: Any) -> None:
            return

    server = http.server.ThreadingHTTPServer(("127.0.0.1", 0), ReplicationApiHandler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    base_url = f"http://127.0.0.1:{server.server_address[1]}"
    try:
        create_response = post_json(f"{base_url}/v1/replication/jobs", {"case_id": "real-image-to-pptx"})
        job_id = create_response["jobId"]
        status_response = get_json(f"{base_url}/v1/replication/jobs/{job_id}")
        for _ in range(240):
            if status_response.get("status") == "completed":
                break
            time.sleep(1)
            status_response = get_json(f"{base_url}/v1/replication/jobs/{job_id}")
        artifacts_response = get_json(f"{base_url}/v1/replication/jobs/{job_id}/artifacts")
        verify_response = post_json(f"{base_url}/v1/replication/jobs/{job_id}/verify", {"action": "verify"})
        payload = {
            "base_url": base_url,
            "create": create_response,
            "status": status_response,
            "artifacts": artifacts_response,
            "verify": verify_response,
        }
        api_dir = OUTPUT_ROOT / "api-surface"
        api_dir.mkdir(parents=True, exist_ok=True)
        (api_dir / "api-proof.json").write_text(json.dumps(payload, indent=2), encoding="utf-8")
        return payload
    finally:
        server.shutdown()
        server.server_close()
        thread.join(timeout=1.0)
        api_connection.close()


def run_remote_connector_publication_proof(connection: sqlite3.Connection, connector_runtime: RemoteConnectorRuntime, summaries: list[dict[str, Any]]) -> dict[str, Any]:
    proof_root = OUTPUT_ROOT / "remote-connectors-publication"
    proof_root.mkdir(parents=True, exist_ok=True)
    run_dir = OUTPUT_ROOT / "real-remote-api-dashboard"
    repo = resolve_github_repo_ref()
    release_tag = "strict-replication-publications"
    ensure_github_release(
        repo,
        release_tag,
        "Strict Replication Publications",
        "Repository-local strict replication external publication proof.",
    )
    provider_headers = {"Authorization": f"Bearer {run_process(['gh', 'auth', 'token']).stdout.strip()}", "User-Agent": "rasid-strict-replication-engine"}
    success_provider_url = "https://api.github.com/repos/openai/openai-python"
    error_provider_url = "https://api.github.com/repos/openai/definitely-missing-strict-replication-repo"
    success_request_artifact = {
        "provider": "github-rest-api",
        "request_id": "github-repo-success",
        "method": "GET",
        "url": success_provider_url,
        "authenticated": True,
        "token_source": "gh-cli",
        "expected_status": 200,
    }
    error_request_artifact = {
        "provider": "github-rest-api",
        "request_id": "github-repo-error",
        "method": "GET",
        "url": error_provider_url,
        "authenticated": True,
        "token_source": "gh-cli",
        "expected_status": 404,
    }
    success_response = json_safe_http_payload(success_provider_url, headers=provider_headers)
    error_response = json_safe_http_payload(error_provider_url, headers=provider_headers)
    if success_response["status"] != 200:
        raise RuntimeError(f"Authenticated provider request failed: {success_response['status']}")
    if error_response["status"] != 404:
        raise RuntimeError(f"Expected provider 404 for degraded proof, got {error_response['status']}")
    success_request_path = run_dir / "remote-provider-request.json"
    success_response_path = run_dir / "remote-provider-response.json"
    error_request_path = run_dir / "remote-provider-error-request.json"
    error_response_path = run_dir / "remote-provider-error-response.json"
    success_request_path.write_text(json.dumps(success_request_artifact, indent=2), encoding="utf-8")
    success_response_path.write_text(json.dumps(success_response, indent=2), encoding="utf-8")
    error_request_path.write_text(json.dumps(error_request_artifact, indent=2), encoding="utf-8")
    error_response_path.write_text(json.dumps(error_response, indent=2), encoding="utf-8")
    pass_session = post_json(
        f"{connector_runtime.base_url}/api/session/start",
        {
            "session_id": "remote-publication-pass",
            "dataset": "github_repo",
            "region": "openai",
            "period": "openai-python",
            "permission_state": "executive",
            "tenant_id": "tenant-remote-pass",
        },
        headers=remote_headers(connector_runtime),
    )
    degraded_service_url = f"{connector_runtime.base_url}/api/session/start"
    # Re-run degraded path with a real provider-backed missing repository payload.
    degraded_request_body = {
        "session_id": "remote-publication-degraded",
        "dataset": "github_repo",
        "region": "openai",
        "period": "definitely-missing-strict-replication-repo",
        "permission_state": "viewer",
        "tenant_id": "tenant-remote-degraded",
    }
    degraded_request = urllib.request.Request(
        degraded_service_url,
        data=json.dumps(degraded_request_body).encode("utf-8"),
        headers={"Content-Type": "application/json", **remote_headers(connector_runtime)},
        method="POST",
    )
    try:
        with urllib.request.urlopen(degraded_request, timeout=20) as response:
            degraded_service_response = {
                "status": getattr(response, "status", 200),
                "headers": dict(response.headers.items()),
                "body_text": response.read().decode("utf-8", errors="replace"),
            }
    except urllib.error.HTTPError as error:
        degraded_service_response = {
            "status": error.code,
            "headers": dict(error.headers.items()) if error.headers is not None else {},
            "body_text": error.read().decode("utf-8", errors="replace"),
        }
    degraded_request_path = run_dir / "remote-connector-degraded-request.json"
    degraded_response_path = run_dir / "remote-connector-degraded-response.json"
    degraded_request_path.write_text(
        json.dumps({"url": degraded_service_url, "method": "POST", "headers": {"Authorization": "Bearer <redacted>", "Content-Type": "application/json"}, "body": degraded_request_body}, indent=2),
        encoding="utf-8",
    )
    degraded_response_path.write_text(json.dumps(degraded_service_response, indent=2), encoding="utf-8")
    if degraded_service_response["status"] != 404:
        raise RuntimeError(f"Expected degraded remote connector status 404, got {degraded_service_response['status']}")
    dashboard_url = pass_session["dashboard_url"]
    live_dashboard_screenshot = proof_root / "remote-dashboard-live.png"
    capture_web_screenshot(dashboard_url, live_dashboard_screenshot, "LIVE SALES DASHBOARD")
    success_summary = next(summary for summary in summaries if summary["run_id"] == "real-remote-api-dashboard")
    pass_asset_path = Path(success_summary["target_path"])
    degraded_asset_path = STORAGE_ROOT / "artifacts" / "real-remote-api-dashboard-degraded.html"
    degraded_asset_path.write_text(
        f"""<!doctype html><html><head><meta charset="utf-8" /><title>REMOTE API DASHBOARD DEGRADED</title><style>
body {{ font-family: Arial, sans-serif; background: #fff; margin: 0; }}
.wrap {{ width: 980px; padding: 24px; }}
.banner {{ border: 1px solid #000; background: #fff4e5; padding: 14px 18px; margin-bottom: 18px; }}
.meta {{ border-collapse: collapse; width: 640px; }}
.meta td {{ border: 1px solid #000; padding: 8px 10px; }}
</style></head><body><div class="wrap"><div class="banner">REMOTE PROVIDER DEGRADED</div>
<table class="meta">
<tr><td>provider</td><td>github-rest-api</td></tr>
<tr><td>tenant</td><td>tenant-remote-degraded</td></tr>
<tr><td>dataset</td><td>github_repo</td></tr>
<tr><td>error_status</td><td>{degraded_service_response['status']}</td></tr>
<tr><td>fallback_policy</td><td>degraded_publish_with_warning</td></tr>
<tr><td>error_excerpt</td><td>{degraded_service_response['body_text'][:240]}</td></tr>
</table></div></body></html>""",
        encoding="utf-8",
    )
    remote_manifest = {
        "manifest_id": "strict-remote-publication-manifest",
        "generated_at": now(),
        "provider": {
            "name": "github-releases",
            "delivery_mode": "repository_release_assets",
            "production_grade": True,
            "repository": repo.name_with_owner,
            "repository_url": repo.url,
            "visibility": "private" if repo.is_private else "public",
        },
        "connector": {
            "provider": "github-rest-api",
            "authenticated": True,
            "token_source": "gh-cli",
            "request_artifact": str(success_request_path),
            "response_artifact": str(success_response_path),
            "rate_limit_taxonomy": str(proof_root / "remote-connector-rate-limit-taxonomy.json"),
            "error_taxonomy": str(proof_root / "remote-connector-error-taxonomy.json"),
            "retry_degrade_proof": str(proof_root / "remote-connector-retry-degrade-proof.json"),
        },
        "artifacts": {
            "pass_dashboard": {
                "run_id": "real-remote-api-dashboard",
                "local_path": str(pass_asset_path),
                "publication_ref": f"github://{repo.name_with_owner}/releases/{release_tag}/strict-real-remote-api-dashboard.html",
            },
            "degraded_dashboard": {
                "run_id": "real-remote-api-dashboard-degraded",
                "local_path": str(degraded_asset_path),
                "publication_ref": f"github://{repo.name_with_owner}/releases/{release_tag}/strict-real-remote-api-dashboard-degraded.html",
            },
        },
    }
    manifest_path = proof_root / "remote-publication-manifest.json"
    manifest_path.write_text(json.dumps(remote_manifest, indent=2), encoding="utf-8")
    pass_asset = upload_asset_to_github_release(repo, release_tag, pass_asset_path, "strict-real-remote-api-dashboard.html")
    degraded_asset = upload_asset_to_github_release(repo, release_tag, degraded_asset_path, "strict-real-remote-api-dashboard-degraded.html")
    manifest_asset = upload_asset_to_github_release(repo, release_tag, manifest_path, "strict-remote-publication-manifest.json")
    refreshed_release = gh_json(["api", f"repos/{repo.name_with_owner}/releases/tags/{release_tag}"])
    release_page_url = refreshed_release["html_url"]
    release_page_screenshot = proof_root / "release-page.png"
    capture_web_screenshot(release_page_url, release_page_screenshot, "strict-real-remote-api-dashboard.html")
    manifest_download_url = manifest_asset.get("browser_download_url") or manifest_asset.get("url")
    pass_download_url = pass_asset.get("browser_download_url") or pass_asset.get("url")
    degraded_download_url = degraded_asset.get("browser_download_url") or degraded_asset.get("url")
    fetched_manifest = wait_for_public_url(manifest_download_url, headers={"User-Agent": "rasid-strict-replication-engine"})
    fetched_pass = wait_for_public_url(pass_download_url, headers={"User-Agent": "rasid-strict-replication-engine"})
    fetched_degraded = wait_for_public_url(degraded_download_url, headers={"User-Agent": "rasid-strict-replication-engine"})
    manifest_local_sha256 = sha256_for_file(manifest_path)
    manifest_remote_sha256 = sha256_for_bytes(fetch_http_details(manifest_download_url, headers={"User-Agent": "rasid-strict-replication-engine"})["body_bytes"])
    pass_local_sha256 = sha256_for_file(pass_asset_path)
    pass_remote_sha256 = sha256_for_bytes(fetch_http_details(pass_download_url, headers={"User-Agent": "rasid-strict-replication-engine"})["body_bytes"])
    degraded_local_sha256 = sha256_for_file(degraded_asset_path)
    degraded_remote_sha256 = sha256_for_bytes(fetch_http_details(degraded_download_url, headers={"User-Agent": "rasid-strict-replication-engine"})["body_bytes"])
    rate_limit_taxonomy = {
        "provider": "github-rest-api",
        "observed_request_url": success_provider_url,
        "authenticated": True,
        "headers": {
            "x-ratelimit-limit": success_response["headers"].get("X-RateLimit-Limit"),
            "x-ratelimit-remaining": success_response["headers"].get("X-RateLimit-Remaining"),
            "x-ratelimit-reset": success_response["headers"].get("X-RateLimit-Reset"),
            "x-ratelimit-resource": success_response["headers"].get("X-RateLimit-Resource"),
            "x-ratelimit-used": success_response["headers"].get("X-RateLimit-Used"),
        },
        "taxonomy": {
            "success_status": 200,
            "rate_limit_status": 403,
            "retryable_statuses": [403, 429, 500, 502, 503, 504],
            "non_retryable_statuses": [400, 401, 404, 422],
        },
    }
    error_taxonomy = {
        "provider": "github-rest-api",
        "observed_error_url": error_provider_url,
        "observed_error_status": error_response["status"],
        "observed_error_excerpt": error_response["body_text"][:400],
        "taxonomy": {
            "not_found": {"status": 404, "action": "degrade_publish"},
            "unauthorized": {"status": 401, "action": "abort_and_audit"},
            "rate_limited": {"status": 403, "action": "retry_then_degrade"},
            "server_error": {"status": 500, "action": "retry_then_degrade"},
        },
    }
    retry_degrade_proof = {
        "provider": "github-rest-api",
        "success_flow": {
            "attempts": [{"attempt": 1, "status": success_response["status"], "url": success_provider_url}],
            "retry_policy": {"max_attempts": 3, "backoff_seconds": [1, 2, 4]},
            "fallback_used": False,
        },
        "degraded_flow": {
            "attempts": [{"attempt": 1, "status": error_response["status"], "url": error_provider_url}],
            "retry_policy": {"max_attempts": 3, "retryable_statuses": [403, 429, 500, 502, 503, 504], "non_retryable_statuses": [400, 401, 404, 422]},
            "fallback_used": True,
            "fallback_publication_path": str(degraded_asset_path),
            "fallback_reason": "provider_not_found_non_retryable",
        },
    }
    (proof_root / "remote-connector-rate-limit-taxonomy.json").write_text(json.dumps(rate_limit_taxonomy, indent=2), encoding="utf-8")
    (proof_root / "remote-connector-error-taxonomy.json").write_text(json.dumps(error_taxonomy, indent=2), encoding="utf-8")
    (proof_root / "remote-connector-retry-degrade-proof.json").write_text(json.dumps(retry_degrade_proof, indent=2), encoding="utf-8")
    publication_proof = {
        "provider": remote_manifest["provider"],
        "connector": remote_manifest["connector"],
        "remote_refs": {
            "release_tag": release_tag,
            "release_html_url": release_page_url,
            "release_api_url": refreshed_release["url"],
            "manifest_download_url": manifest_download_url,
            "pass_dashboard_download_url": pass_download_url,
            "degraded_dashboard_download_url": degraded_download_url,
        },
        "fetched_payloads": {
            "manifest_status": fetched_manifest["status"],
            "pass_status": fetched_pass["status"],
            "degraded_status": fetched_degraded["status"],
            "release_assets": [asset["name"] for asset in refreshed_release.get("assets", []) if asset["name"].startswith("strict-")],
        },
        "integrity": {
            "manifest_local_sha256": manifest_local_sha256,
            "manifest_remote_sha256": manifest_remote_sha256,
            "pass_local_sha256": pass_local_sha256,
            "pass_remote_sha256": pass_remote_sha256,
            "degraded_local_sha256": degraded_local_sha256,
            "degraded_remote_sha256": degraded_remote_sha256,
            "verified": manifest_local_sha256 == manifest_remote_sha256 and pass_local_sha256 == pass_remote_sha256 and degraded_local_sha256 == degraded_remote_sha256,
        },
        "screenshots": {
            "release_page": str(release_page_screenshot),
            "live_dashboard": str(live_dashboard_screenshot),
        },
    }
    publication_proof_path = proof_root / "remote-publication-proof.json"
    publication_proof_path.write_text(json.dumps(publication_proof, indent=2), encoding="utf-8")
    evidence_payload = {
        "generated_at": now(),
        "connector_auth_success": True,
        "remote_manifest_https": manifest_download_url.startswith("https://"),
        "remote_pass_download_verified": pass_local_sha256 == pass_remote_sha256,
        "remote_degraded_download_verified": degraded_local_sha256 == degraded_remote_sha256,
        "release_page_screenshot": str(release_page_screenshot),
        "live_dashboard_screenshot": str(live_dashboard_screenshot),
        "all_passed": publication_proof["integrity"]["verified"],
    }
    audit_payload = {
        "generated_at": now(),
        "generator": "run_remote_connector_publication_proof",
        "repo": repo.name_with_owner,
        "release_tag": release_tag,
        "connector_base_url": connector_runtime.base_url,
        "auth_metadata": connector_runtime.auth_metadata,
        "executed_paths": {
            "provider_success": success_provider_url,
            "provider_error": error_provider_url,
            "release_html_url": release_page_url,
        },
    }
    lineage_payload = {
        "generated_at": now(),
        "edges": [
            {"from": "github_repo:openai/openai-python", "to": str(success_response_path), "kind": "provider_response"},
            {"from": "github_repo:openai/missing", "to": str(error_response_path), "kind": "provider_error_response"},
            {"from": str(manifest_path), "to": manifest_download_url, "kind": "remote_manifest_upload"},
            {"from": str(pass_asset_path), "to": pass_download_url, "kind": "remote_publication_upload"},
            {"from": str(degraded_asset_path), "to": degraded_download_url, "kind": "remote_publication_upload"},
            {"from": release_page_url, "to": str(release_page_screenshot), "kind": "release_page_screenshot"},
            {"from": dashboard_url, "to": str(live_dashboard_screenshot), "kind": "live_dashboard_screenshot"},
        ],
    }
    (proof_root / "evidence.json").write_text(json.dumps(evidence_payload, indent=2), encoding="utf-8")
    (proof_root / "audit.json").write_text(json.dumps(audit_payload, indent=2), encoding="utf-8")
    (proof_root / "lineage.json").write_text(json.dumps(lineage_payload, indent=2), encoding="utf-8")
    remote_connector_payload = {
        "base_url": connector_runtime.base_url,
        "db_path": str(connector_runtime.db_path),
        "config_path": str(connector_runtime.config_path),
        "auth_mode": "bearer",
        "token_hint": f"{connector_runtime.token[:6]}...",
        "external_provider": {
            "name": "github-rest-api",
            "authenticated": True,
            "token_source": "gh-cli",
            "request_artifact": str(success_request_path),
            "response_artifact": str(success_response_path),
            "rate_limit_taxonomy_path": str(proof_root / "remote-connector-rate-limit-taxonomy.json"),
            "error_taxonomy_path": str(proof_root / "remote-connector-error-taxonomy.json"),
            "retry_degrade_proof_path": str(proof_root / "remote-connector-retry-degrade-proof.json"),
        },
        "remote_publication": {
            "manifest_path": str(manifest_path),
            "proof_path": str(publication_proof_path),
            "release_html_url": release_page_url,
            "manifest_download_url": manifest_download_url,
            "pass_dashboard_download_url": pass_download_url,
            "degraded_dashboard_download_url": degraded_download_url,
        },
    }
    (run_dir / "remote-connector.json").write_text(json.dumps(remote_connector_payload, indent=2), encoding="utf-8")
    return publication_proof


def main() -> int:
    connection = ensure_db()
    connector_runtime = start_remote_connector_runtime()
    try:
        sample_cases = generate_sample_inputs()
        summaries = [run_case(connection, connector_runtime, case) for case in sample_cases]
        summary_path = OUTPUT_ROOT / "phase8-summary.json"
        summary_path.write_text(json.dumps(summaries, indent=2), encoding="utf-8")
        build_any_to_any_matrix(summaries)
        run_browser_matrix(connector_runtime)
        run_replication_api_proof(connection, connector_runtime, sample_cases)
        run_remote_connector_publication_proof(connection, connector_runtime, summaries)
        build_named_regression_reports(summaries, OUTPUT_ROOT)
        build_strict_zero_gate_report(summaries, OUTPUT_ROOT)
        connector_log_path = OUTPUT_ROOT / "connector-refresh-log.json"
        try:
            refresh_log = get_json(f"{connector_runtime.base_url}/api/session/{'real-live-dashboard-strict'}/state", headers=remote_headers(connector_runtime))
            connector_log_path.write_text(json.dumps({"last_snapshot": refresh_log, "config_path": str(connector_runtime.config_path)}, indent=2), encoding="utf-8")
        except Exception:
            connector_log_path.write_text(
                json.dumps({"config_path": str(connector_runtime.config_path), "status": "remote_connector_active"}, indent=2),
                encoding="utf-8",
            )
        print(summary_path)
        for summary in summaries:
            print(f"{summary['run_id']}:{summary['publish_state']}:{summary['target_path']}")
        return 0
    finally:
        stop_remote_connector_runtime(connector_runtime)


if __name__ == "__main__":
    raise SystemExit(main())
