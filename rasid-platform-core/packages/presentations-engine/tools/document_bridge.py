import argparse
import base64
import io
import json
import os
import re
import unicodedata
import zipfile
from pathlib import Path

import fitz
from openpyxl import load_workbook
from pptx import Presentation
from reportlab.lib.colors import HexColor
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from bidi.algorithm import get_display
import arabic_reshaper


def emu_to_px(value):
    try:
        return max(1, round(int(value) / 9525))
    except Exception:
        return 1


def ensure_font(name, font_path):
    try:
        pdfmetrics.getFont(name)
    except Exception:
        pdfmetrics.registerFont(TTFont(name, font_path))


def shape_text(value, rtl):
    if not value:
        return ""
    text = str(value)
    return get_display(arabic_reshaper.reshape(text)) if rtl else text


def normalize_semantic_text(value):
    if not value:
        return ""
    text = unicodedata.normalize("NFKC", str(value))
    text = re.sub(r"[\u200e\u200f\u202a-\u202e\u2066-\u2069]", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def extract_theme_tokens(archive):
    try:
        xml = archive.read("ppt/theme/theme1.xml").decode("utf-8", "ignore")
    except KeyError:
        return {
            "primary_color": "C7511F",
            "secondary_color": "0F172A",
            "accent_color": "1D8F6E",
            "neutral_color": "EEF2F6",
            "font_face": "Aptos",
        }
    colors = re.findall(r'<a:srgbClr val="([A-Fa-f0-9]{6})"', xml)
    font_face = re.search(r'<a:latin typeface="([^"]+)"', xml)
    return {
        "primary_color": colors[0].upper() if len(colors) > 0 else "C7511F",
        "secondary_color": colors[1].upper() if len(colors) > 1 else "0F172A",
        "accent_color": colors[2].upper() if len(colors) > 2 else "1D8F6E",
        "neutral_color": colors[5].upper() if len(colors) > 5 else "EEF2F6",
        "font_face": font_face.group(1) if font_face else "Aptos",
    }


def extract_pptx(input_path):
    prs = Presentation(input_path)
    archive = zipfile.ZipFile(input_path)
    slide_width_px = emu_to_px(prs.slide_width)
    slide_height_px = emu_to_px(prs.slide_height)
    theme_tokens = extract_theme_tokens(archive)
    slide_titles = []
    layout_refs = []
    slide_geometries = []
    media_refs = []
    notes = []
    text_parts = []
    layouts = []
    masters = []

    for slide_index, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name or f"layout-{slide_index + 1}"
        master_name = slide.slide_layout.slide_master.name or "master-default"
        layout_refs.append(layout_name)
        layouts.append(layout_name)
        masters.append(master_name)

        notes_text = ""
        try:
            notes_text = slide.notes_slide.notes_text_frame.text or ""
        except Exception:
            notes_text = ""
        if notes_text:
            notes.append(notes_text)

        shapes = []
        title_text = ""
        for shape in slide.shapes:
            entry = {
                "x": emu_to_px(getattr(shape, "left", 0)),
                "y": emu_to_px(getattr(shape, "top", 0)),
                "w": emu_to_px(getattr(shape, "width", 0)),
                "h": emu_to_px(getattr(shape, "height", 0)),
                "shape_name": getattr(shape, "name", ""),
                "element_type": str(getattr(shape, "shape_type", "")),
                "placeholder_type": None,
                "placeholder_idx": None,
                "layout_ref": layout_name,
                "master_ref": master_name,
                "text": "",
            }
            if getattr(shape, "is_placeholder", False):
                try:
                    entry["placeholder_type"] = str(shape.placeholder_format.type)
                    entry["placeholder_idx"] = shape.placeholder_format.idx
                except Exception:
                    pass
            if getattr(shape, "has_text_frame", False):
                entry["text"] = shape.text or ""
                if not title_text and entry["text"]:
                    if "title" in (entry["placeholder_type"] or "").lower() or "title" in entry["shape_name"].lower():
                        title_text = entry["text"]
                if entry["text"]:
                    text_parts.append(entry["text"])
            if hasattr(shape, "image"):
                try:
                    entry["image_ext"] = shape.image.ext
                except Exception:
                    pass
            shapes.append(entry)

        if not title_text:
            try:
                title_text = slide.shapes.title.text
            except Exception:
                title_text = ""
        if not title_text:
            text_shape = next((item["text"] for item in shapes if item["text"]), "")
            title_text = text_shape or f"Slide {slide_index + 1}"
        slide_titles.append(title_text)

        rel_targets = []
        try:
            for rel in slide.part.rels.values():
                target = str(rel.target_ref)
                rel_targets.append(target)
                if "media/" in target:
                    media_refs.append(target)
        except Exception:
            pass

        slide_geometries.append(
            {
                "slide_index": slide_index,
                "width_px": slide_width_px,
                "height_px": slide_height_px,
                "layout_ref": layout_name,
                "master_ref": master_name,
                "notes_text": notes_text,
                "media_refs": rel_targets,
                "element_boxes": shapes,
            }
        )

    return {
        "page_count": len(prs.slides),
        "slide_titles": slide_titles,
        "layout_refs": layout_refs,
        "master_refs": masters,
        "notes_count": len(notes),
        "notes_text": notes,
        "media_refs": sorted(set(media_refs)),
        "theme_tokens": theme_tokens,
        "slide_geometries": slide_geometries,
        "extracted_text": "\n".join(text_parts),
    }


def extract_pdf(input_path, render_dir=None):
    doc = fitz.open(input_path)
    pages = []
    texts = []
    rendered_pages = []
    total_images = 0
    total_tables = 0
    for page_index, page in enumerate(doc):
        text = page.get_text("text", sort=True)
        texts.append(text)
        normalized_text = normalize_semantic_text(text)
        images = page.get_images(full=True)
        total_images += len(images)
        table_count = 0
        try:
            tables = page.find_tables()
            table_count = len(getattr(tables, "tables", []))
        except Exception:
            table_count = 0
        total_tables += table_count
        if render_dir:
            Path(render_dir).mkdir(parents=True, exist_ok=True)
            png_path = Path(render_dir) / f"page-{page_index + 1}.png"
            page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False).save(str(png_path))
            rendered_pages.append(str(png_path))
        pages.append(
            {
                "page_index": page_index,
                "width": page.rect.width,
                "height": page.rect.height,
                "text": text,
                "normalized_text": normalized_text,
                "image_count": len(images),
                "table_count": table_count,
            }
        )
    return {
        "page_count": doc.page_count,
        "pages": pages,
        "extracted_text": "\n".join(texts),
        "normalized_text": "\n".join(
            page["normalized_text"] for page in pages if page.get("normalized_text")
        ),
        "image_count": total_images,
        "table_count": total_tables,
        "rendered_pages": rendered_pages,
    }


def extract_xlsx(input_path):
    workbook = load_workbook(input_path, data_only=True, read_only=True)
    sheets = []
    extracted_lines = []
    dataset_summaries = []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(value) if value is not None else f"Column{index + 1}" for index, value in enumerate(rows[0])]
        sample_rows = []
        for row in rows[1:21]:
            mapped = {}
            has_value = False
            for index, header in enumerate(headers):
                value = row[index] if index < len(row) else None
                mapped[header] = value
                has_value = has_value or value not in (None, "")
            if has_value:
                sample_rows.append(mapped)
        sheets.append(
            {
                "sheet_name": sheet.title,
                "headers": headers,
                "row_count": max(0, sheet.max_row - 1),
                "column_count": sheet.max_column,
                "sample_rows": sample_rows[:10],
            }
        )
        dataset_summaries.append(
            {
                "sheet_name": sheet.title,
                "headers": headers,
                "rows": sample_rows[:10],
            }
        )
        extracted_lines.append(sheet.title)
        for row in rows[:5]:
            extracted_lines.append(" | ".join("" if value is None else str(value) for value in row))
    return {
        "sheet_count": len(sheets),
        "sheet_names": [sheet["sheet_name"] for sheet in sheets],
        "sheets": sheets,
        "dataset_summaries": dataset_summaries,
        "extracted_text": "\n".join(line for line in extracted_lines if line),
    }


def decode_media(spec):
    if spec.get("file_path") and Path(spec["file_path"]).exists():
        return ImageReader(spec["file_path"])
    data_uri = spec.get("data_uri")
    if data_uri and data_uri.startswith("data:"):
        encoded = data_uri.split(",", 1)[1]
        return ImageReader(io.BytesIO(base64.b64decode(encoded)))
    return None


def render_pdf(spec_path, output_path):
    spec = json.loads(Path(spec_path).read_text(encoding="utf-8"))
    font_regular = r"C:\Windows\Fonts\arial.ttf"
    font_bold = r"C:\Windows\Fonts\arialbd.ttf"
    ensure_font("ArialArabic", font_regular)
    ensure_font("ArialArabicBold", font_bold if Path(font_bold).exists() else font_regular)

    page_width = 960
    page_height = 540
    deck = spec["deck"]
    pdf = canvas.Canvas(output_path, pagesize=(page_width, page_height))
    pdf.setTitle(deck.get("title", "Presentation"))
    pdf.setAuthor("Rasid Presentations Engine")
    rtl = deck.get("rtl", False)

    for slide in spec["slides"]:
        pdf.setFillColor(HexColor("#FFFFFF"))
        pdf.rect(0, 0, page_width, page_height, stroke=0, fill=1)
        pdf.setFillColor(HexColor("#0F172A"))
        for block in slide["blocks"]:
            x = block["box"]["x"] * 72 / 96
            y = page_height - ((block["box"]["y"] + block["box"]["h"]) * 72 / 96)
            w = block["box"]["w"] * 72 / 96
            h = block["box"]["h"] * 72 / 96
            kind = block["kind"]
            title = shape_text(block.get("title", ""), rtl)
            body = shape_text(block.get("body", ""), rtl)
            if kind in {"image", "video"} and block.get("media"):
                image = decode_media(block["media"])
                if image:
                    try:
                        pdf.drawImage(image, x, y, width=w, height=h, preserveAspectRatio=True, mask="auto")
                    except Exception:
                        pass
                if block["media"].get("caption"):
                    pdf.setFont("ArialArabic", 9)
                    pdf.setFillColor(HexColor("#334155"))
                    pdf.drawString(x, max(12, y - 10), shape_text(block["media"]["caption"], rtl))
                continue
            if kind == "table" and block.get("table"):
                table = block["table"]
                pdf.setFont("ArialArabicBold", 10)
                cursor_y = y + h - 14
                pdf.drawString(x, cursor_y, shape_text(" | ".join(table["columns"]), rtl))
                pdf.setFont("ArialArabic", 9)
                for row in table["rows"][:8]:
                    cursor_y -= 14
                    pdf.drawString(x, cursor_y, shape_text(" | ".join(row), rtl))
                continue
            if kind == "chart" and block.get("chart"):
                chart = block["chart"]
                values = chart["series"][0]["values"] if chart["series"] else [1]
                max_value = max(values) if values else 1
                for index, value in enumerate(values[:6]):
                    bar_y = y + h - 20 - index * 18
                    width = max(12, ((w - 120) * value / max_value)) if max_value else 12
                    pdf.setFillColor(HexColor("#C7511F"))
                    pdf.rect(x + 90, bar_y, width, 10, stroke=0, fill=1)
                    pdf.setFillColor(HexColor("#0F172A"))
                    pdf.setFont("ArialArabic", 8)
                    label = chart["categories"][index] if index < len(chart["categories"]) else f"Row {index + 1}"
                    pdf.drawString(x, bar_y + 1, shape_text(str(label), rtl))
                continue
            pdf.setFillColor(HexColor("#0F172A"))
            pdf.setFont("ArialArabicBold" if kind == "title" else "ArialArabic", 18 if kind == "title" else 10)
            if title:
                pdf.drawString(x, y + h - 18, title)
            if body:
                pdf.setFont("ArialArabic", 10)
                cursor_y = y + h - 34
                for line in body.split("\n")[:10]:
                    pdf.drawString(x, cursor_y, line)
                    cursor_y -= 12
        pdf.showPage()
    pdf.save()


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("command")
    parser.add_argument("input")
    parser.add_argument("output")
    parser.add_argument("--render-dir", default=None)
    parser.add_argument("--pdf-path", default=None)
    args = parser.parse_args()

    if args.command == "extract-pptx":
        payload = extract_pptx(args.input)
    elif args.command == "extract-pdf":
        payload = extract_pdf(args.input, render_dir=args.render_dir)
    elif args.command == "extract-xlsx":
        payload = extract_xlsx(args.input)
    elif args.command == "render-pdf":
        if not args.pdf_path:
            raise SystemExit("--pdf-path is required for render-pdf")
        render_pdf(args.input, args.pdf_path)
        payload = {"output_path": args.pdf_path}
    else:
        raise SystemExit(f"Unknown command: {args.command}")

    Path(args.output).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


if __name__ == "__main__":
    main()
